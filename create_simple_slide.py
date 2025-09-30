#!/usr/bin/env python3
"""
Direct script to create a simple presentation and upload to S3
"""

import boto3
import json
import base64
import io
import tempfile
import os
from datetime import datetime

# Import only what we need
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

def create_basic_json_data():
    """Create the loan portfolio data as JSON that can be uploaded to S3"""
    data = {
        "title": "Loan Portfolio",
        "subtitle": "Total Loans Held for Investment ($ in Millions)",
        "data": {
            "quarters": ["2Q'19", "3Q'19", "4Q'19", "1Q'20", "2Q'20"],
            "loan_amounts": [1936, 1963, 2144, 2109, 2332],
            "yield_percentages": [5.90, 5.91, 5.79, 5.76, 5.26],
            "ppp_yield": 5.06
        },
        "highlights": {
            "title": "2Q'20 Highlights",
            "items": [
                "Total loan increase of $229.9M vs. 1Q'20",
                "Growth from $215.3M PPP loans and $34.7M seasonal agriculture loans",
                "Partial offset from $24.4M pay-downs in non-residential consumer and direct energy loans",
                "Over 2,000 PPP loans closed",
                "2Q'20 yield of 5.26% (down 50 bps vs. 1Q'20 excluding PPP)"
            ]
        },
        "styling": {
            "primary_color": "#C00000",
            "secondary_color": "#808080",
            "footer_text": "South Plains Financial, Inc."
        },
        "generated_at": datetime.now().isoformat()
    }
    return data

def upload_json_to_s3(data, filename):
    """Upload JSON data to S3"""
    s3 = boto3.client('s3')
    bucket = 'scribbe-ai-dev-output'
    
    # Convert to JSON
    json_content = json.dumps(data, indent=2)
    
    # Upload to S3
    s3.put_object(
        Bucket=bucket,
        Key=filename,
        Body=json_content,
        ContentType='application/json'
    )
    
    return f"s3://{bucket}/{filename}"

def create_pptx_if_possible():
    """Create a PowerPoint if python-pptx is available"""
    if not HAS_PPTX:
        return None
        
    try:
        # Create presentation
        prs = Presentation()
        
        # Use blank layout
        slide_layout = prs.slide_layouts[5]  # Blank slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Loan Portfolio"
        title_frame.paragraphs[0].font.name = "Arial"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(5), Inches(0.4))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Total Loans Held for Investment ($ in Millions)"
        subtitle_frame.paragraphs[0].font.size = Pt(14)
        subtitle_frame.paragraphs[0].font.bold = True
        
        # Add data table
        quarters = ["2Q'19", "3Q'19", "4Q'19", "1Q'20", "2Q'20"]
        amounts = [1936, 1963, 2144, 2109, 2332]
        yields = [5.90, 5.91, 5.79, 5.76, 5.26]
        
        # Create data display
        data_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(3))
        data_frame = data_box.text_frame
        
        # Add headers
        data_frame.text = f"{'Quarter':<10} {'Loans ($M)':<12} {'Yield (%)':<10}"
        data_frame.paragraphs[0].font.bold = True
        data_frame.paragraphs[0].font.size = Pt(12)
        
        # Add data rows
        for i, (quarter, amount, yield_pct) in enumerate(zip(quarters, amounts, yields)):
            p = data_frame.add_paragraph()
            p.text = f"{quarter:<10} {amount:<12} {yield_pct:<10}"
            p.font.size = Pt(11)
            if quarter == "2Q'20":
                p.font.color.rgb = RGBColor(192, 0, 0)  # Red for latest quarter
        
        # Add PPP yield note
        p = data_frame.add_paragraph()
        p.text = f"{'2Q\'20 PPP':<10} {'':<12} {5.06:<10}"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(128, 128, 128)  # Gray for PPP
        
        # Add highlights section
        highlights_box = slide.shapes.add_textbox(Inches(6.3), Inches(1.2), Inches(3.5), Inches(3.8))
        highlights_frame = highlights_box.text_frame
        highlights_frame.word_wrap = True
        
        # Add highlights background
        highlights_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.2), Inches(1.1),
            Inches(3.7), Inches(4)
        )
        highlights_bg.fill.solid()
        highlights_bg.fill.fore_color.rgb = RGBColor(245, 245, 245)
        highlights_bg.line.fill.background()
        
        # Move highlights text to front
        slide.shapes._spTree.remove(highlights_box._element)
        slide.shapes._spTree.append(highlights_box._element)
        
        # Add highlights title
        highlights_frame.text = "2Q'20 Highlights"
        highlights_frame.paragraphs[0].font.size = Pt(14)
        highlights_frame.paragraphs[0].font.bold = True
        highlights_frame.paragraphs[0].font.color.rgb = RGBColor(192, 0, 0)
        
        # Add highlight items
        highlights = [
            "• Total loan increase of $229.9M vs. 1Q'20",
            "• Growth from $215.3M PPP loans and $34.7M seasonal agriculture loans", 
            "• Partial offset from $24.4M pay-downs in non-residential consumer and direct energy loans",
            "• Over 2,000 PPP loans closed",
            "• 2Q'20 yield of 5.26% (down 50 bps vs. 1Q'20 excluding PPP)"
        ]
        
        for highlight in highlights:
            p = highlights_frame.add_paragraph()
            p.text = highlight
            p.font.size = Pt(10)
            p.space_after = Pt(6)
        
        # Add footer
        footer_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(6.8),
            Inches(10), Inches(0.4)
        )
        footer_shape.fill.solid()
        footer_shape.fill.fore_color.rgb = RGBColor(128, 128, 128)
        footer_shape.line.fill.background()
        
        footer_text = slide.shapes.add_textbox(
            Inches(0.3), Inches(6.85),
            Inches(9.4), Inches(0.3)
        )
        footer_frame = footer_text.text_frame
        footer_frame.text = "South Plains Financial, Inc."
        footer_frame.paragraphs[0].font.size = Pt(12)
        footer_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        return prs
        
    except Exception as e:
        print(f"Error creating PowerPoint: {e}")
        return None

def main():
    """Main function"""
    print("🚀 Creating Loan Portfolio slide data...")
    
    # Create the data
    loan_data = create_basic_json_data()
    
    # Upload JSON data
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    json_filename = f"loan_portfolio/data_{timestamp}.json"
    
    try:
        json_url = upload_json_to_s3(loan_data, json_filename)
        print(f"✅ JSON data uploaded to: {json_url}")
        
        # Try to create PowerPoint if possible
        prs = create_pptx_if_possible()
        if prs:
            # Save to temporary file
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                prs.save(tmp.name)
                
                # Upload PowerPoint to S3
                s3 = boto3.client('s3')
                pptx_filename = f"loan_portfolio/presentation_{timestamp}.pptx"
                
                with open(tmp.name, 'rb') as f:
                    s3.put_object(
                        Bucket='scribbe-ai-dev-output',
                        Key=pptx_filename,
                        Body=f.read(),
                        ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                    )
                
                # Clean up
                os.unlink(tmp.name)
                
                print(f"✅ PowerPoint uploaded to: s3://scribbe-ai-dev-output/{pptx_filename}")
                
                # Generate presigned URL
                try:
                    presigned_url = s3.generate_presigned_url(
                        'get_object',
                        Params={'Bucket': 'scribbe-ai-dev-output', 'Key': pptx_filename},
                        ExpiresIn=3600
                    )
                    print(f"\n📥 Download PowerPoint (valid for 1 hour):")
                    print(f"   {presigned_url}")
                except Exception as e:
                    print(f"Could not generate presigned URL: {e}")
        else:
            print("⚠️  Could not create PowerPoint (python-pptx not available)")
            
    except Exception as e:
        print(f"❌ Error: {str(e)}")

if __name__ == "__main__":
    main()