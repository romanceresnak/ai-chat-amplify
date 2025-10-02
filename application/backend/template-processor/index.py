import json
import boto3
import os
import logging
from typing import Dict, Any, List
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
from datetime import datetime

# Initialize logging
logger = logging.getLogger()
logger.setLevel(logging.INFO)

# Initialize AWS clients
s3 = boto3.client('s3')

# Environment variables
ENVIRONMENT = os.environ['ENVIRONMENT']
TEMPLATES_BUCKET = os.environ['TEMPLATES_BUCKET']
OUTPUT_BUCKET = os.environ['OUTPUT_BUCKET']

def lambda_handler(event: Dict[str, Any], context: Any) -> Dict[str, Any]:
    """
    Process PowerPoint template and generate final presentation.
    """
    try:
        logger.info(f"Received event: {json.dumps(event)}")
        
        # Parse request
        template_key = event.get('template_key')
        presentation_content = event.get('presentation_content', {})
        output_key = event.get('output_key')
        
        if not template_key or not output_key:
            return {
                'statusCode': 400,
                'body': json.dumps({
                    'error': 'Missing required parameters: template_key and output_key'
                })
            }
        
        # Download template from S3
        template_path = download_template(template_key)
        
        # Open the presentation
        prs = Presentation(template_path)
        
        # Process slides
        process_presentation(prs, presentation_content)
        
        # Save presentation
        output_buffer = io.BytesIO()
        prs.save(output_buffer)
        output_buffer.seek(0)
        
        # Upload to S3
        upload_presentation(output_buffer, output_key)
        
        # Clean up
        os.remove(template_path)
        
        return {
            'statusCode': 200,
            'body': json.dumps({
                'message': 'Presentation processed successfully',
                'output_key': output_key
            })
        }
        
    except Exception as e:
        logger.error(f"Error in template processor: {str(e)}", exc_info=True)
        return {
            'statusCode': 500,
            'body': json.dumps({
                'error': 'Internal server error',
                'message': str(e)
            })
        }

def download_template(template_key: str) -> str:
    """
    Download template from S3.
    """
    try:
        local_path = f"/tmp/template_{datetime.now().timestamp()}.pptx"
        
        # Handle different template_key formats
        if template_key.startswith('templates/') and template_key.endswith('.pptx'):
            # Frontend sends "templates/default.pptx", but file is at "default.pptx"
            template_name = template_key.replace('templates/', '')
            s3_key = template_name
        elif template_key.endswith('.pptx'):
            # Direct file path like "default.pptx"
            s3_key = template_key
        else:
            # Directory path like "default"
            s3_key = f"{template_key}.pptx"
        
        logger.info(f"Downloading template from s3://{TEMPLATES_BUCKET}/{s3_key}")
        s3.download_file(
            Bucket=TEMPLATES_BUCKET,
            Key=s3_key,
            Filename=local_path
        )
        
        return local_path
        
    except Exception as e:
        logger.error(f"Error downloading template: {str(e)}")
        raise

def process_presentation(prs: Presentation, content: Dict[str, Any]):
    """
    Process the presentation with generated content.
    """
    try:
        slides_content = content.get('slides', [])
        metadata = content.get('metadata', {})
        
        # Update presentation metadata
        update_presentation_metadata(prs, metadata)
        
        # Process each slide
        for slide_content in slides_content:
            slide_number = slide_content.get('slide_number', 1)
            slide_type = slide_content.get('slide_type', 'content')
            
            # Get or create slide
            if slide_number <= len(prs.slides):
                slide = prs.slides[slide_number - 1]
            else:
                # Add new slide with appropriate layout
                layout = get_slide_layout(prs, slide_type)
                slide = prs.slides.add_slide(layout)
            
            # Process slide based on type
            if slide_type == 'title':
                process_title_slide(slide, slide_content)
            elif slide_type == 'executive_summary':
                process_executive_summary_slide(slide, slide_content)
            elif slide_type == 'financial_overview':
                process_financial_slide(slide, slide_content)
            elif slide_type == 'chart':
                process_chart_slide(slide, slide_content)
            elif slide_type == 'table':
                process_table_slide(slide, slide_content)
            elif slide_type == 'content':
                process_content_slide(slide, slide_content)
            else:
                process_generic_slide(slide, slide_content)
            
            # Add speaker notes
            if slide_content.get('notes'):
                slide.notes_slide.notes_text_frame.text = slide_content['notes']
                
    except Exception as e:
        logger.error(f"Error processing presentation: {str(e)}")
        raise

def update_presentation_metadata(prs: Presentation, metadata: Dict[str, Any]):
    """
    Update presentation core properties.
    """
    try:
        core_props = prs.core_properties
        core_props.title = metadata.get('title', 'ScribbeAI Presentation')
        core_props.author = 'ScribbeAI'
        core_props.comments = metadata.get('description', '')
        core_props.created = datetime.now()
        
    except Exception as e:
        logger.warning(f"Could not update presentation metadata: {str(e)}")

def get_slide_layout(prs: Presentation, slide_type: str):
    """
    Get appropriate slide layout based on slide type.
    """
    layout_mapping = {
        'title': 0,  # Title slide
        'content': 1,  # Title and content
        'section': 2,  # Section header
        'two_content': 3,  # Two content
        'comparison': 4,  # Comparison
        'title_only': 5,  # Title only
        'blank': 6  # Blank
    }
    
    # Map slide types to layouts
    type_to_layout = {
        'title': 'title',
        'executive_summary': 'content',
        'financial_overview': 'content',
        'chart': 'content',
        'table': 'content',
        'content': 'content',
        'section': 'section'
    }
    
    layout_name = type_to_layout.get(slide_type, 'content')
    layout_index = layout_mapping.get(layout_name, 1)
    
    return prs.slide_layouts[layout_index]

def process_title_slide(slide, content: Dict[str, Any]):
    """
    Process title slide with company info and date.
    """
    try:
        # Title
        if slide.shapes.title:
            slide.shapes.title.text = content.get('title', '')
        
        # Subtitle (usually contains company name and date)
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Subtitle placeholder
                shape.text = content.get('content', {}).get('subtitle', '')
                break
                
    except Exception as e:
        logger.error(f"Error processing title slide: {str(e)}")

def process_executive_summary_slide(slide, content: Dict[str, Any]):
    """
    Process executive summary slide with key highlights.
    """
    try:
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = content.get('title', 'Executive Summary')
        
        # Add content
        content_data = content.get('content', {})
        
        # Find content placeholder
        for shape in slide.placeholders:
            if hasattr(shape, 'text_frame') and shape.placeholder_format.idx > 0:
                tf = shape.text_frame
                tf.clear()
                
                # Add highlights
                highlights = content_data.get('highlights', [])
                for highlight in highlights:
                    p = tf.add_paragraph()
                    p.text = f"• {highlight}"
                    p.level = 0
                
                # Add key metrics
                if content_data.get('key_metrics'):
                    p = tf.add_paragraph()
                    p.text = "\nKey Metrics:"
                    p.level = 0
                    
                    for metric, value in content_data.get('key_metrics', {}).items():
                        p = tf.add_paragraph()
                        p.text = f"• {metric}: {value}"
                        p.level = 1
                
                break
                
    except Exception as e:
        logger.error(f"Error processing executive summary slide: {str(e)}")

def process_financial_slide(slide, content: Dict[str, Any]):
    """
    Process financial overview slides with metrics and trends.
    """
    try:
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = content.get('title', 'Financial Overview')
        
        content_data = content.get('content', {})
        
        # Process charts if any
        charts_data = content.get('charts', [])
        if charts_data:
            add_charts_to_slide(slide, charts_data)
        
        # Process tables if any
        tables_data = content.get('tables', [])
        if tables_data:
            add_tables_to_slide(slide, tables_data)
        
        # Add text content
        for shape in slide.placeholders:
            if hasattr(shape, 'text_frame') and shape.placeholder_format.idx > 0:
                tf = shape.text_frame
                tf.clear()
                
                # Add financial metrics
                if isinstance(content_data, dict):
                    for section, data in content_data.items():
                        p = tf.add_paragraph()
                        p.text = section.replace('_', ' ').title()
                        p.level = 0
                        
                        if isinstance(data, dict):
                            for key, value in data.items():
                                p = tf.add_paragraph()
                                p.text = f"• {key}: {value}"
                                p.level = 1
                
                break
                
    except Exception as e:
        logger.error(f"Error processing financial slide: {str(e)}")

def process_chart_slide(slide, content: Dict[str, Any]):
    """
    Process slides with charts and visualizations.
    """
    try:
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = content.get('title', '')
        
        # Check if this slide needs highlights layout
        if content.get('highlights'):
            # Create slide with highlights panel
            chart_position = create_slide_with_highlights_layout(slide, content)
            if chart_position:
                # Add chart to specific position
                charts_data = content.get('charts', [])
                if charts_data:
                    chart_spec = charts_data[0]  # Use first chart for highlights layout
                    chart_type = chart_spec.get('chart_type', 'bar')
                    x, y, cx, cy = chart_position
                    
                    # Add the chart
                    if chart_type == 'bar':
                        add_bar_chart(slide, chart_spec, x, y, cx, cy)
                    elif chart_type == 'line':
                        add_line_chart(slide, chart_spec, x, y, cx, cy)
                    elif chart_type == 'pie':
                        add_pie_chart(slide, chart_spec, x, y, cx, cy)
                    elif chart_type == 'donut':
                        add_donut_chart(slide, chart_spec, x, y, cx, cy)
                    elif chart_type == 'combo':
                        add_combo_chart(slide, chart_spec, x, y, cx, cy)
        else:
            # Normal chart layout
            charts_data = content.get('charts', [])
            add_charts_to_slide(slide, charts_data)
        
        # Apply corporate styling if specified
        if content.get('style'):
            apply_corporate_styling(slide, content.get('style'))
        
    except Exception as e:
        logger.error(f"Error processing chart slide: {str(e)}")

def add_charts_to_slide(slide, charts_data: List[Dict[str, Any]]):
    """
    Add charts to slide based on specifications.
    """
    try:
        for i, chart_spec in enumerate(charts_data):
            chart_type = chart_spec.get('chart_type', 'bar')
            
            # Define chart position and size
            x, y, cx, cy = calculate_chart_position(i, len(charts_data))
            
            # Create chart based on type
            if chart_type == 'bar':
                add_bar_chart(slide, chart_spec, x, y, cx, cy)
            elif chart_type == 'line':
                add_line_chart(slide, chart_spec, x, y, cx, cy)
            elif chart_type == 'pie':
                add_pie_chart(slide, chart_spec, x, y, cx, cy)
            elif chart_type == 'donut':
                add_donut_chart(slide, chart_spec, x, y, cx, cy)
            elif chart_type == 'combo':
                add_combo_chart(slide, chart_spec, x, y, cx, cy)
            else:
                logger.warning(f"Unsupported chart type: {chart_type}")
                
    except Exception as e:
        logger.error(f"Error adding charts to slide: {str(e)}")

def calculate_chart_position(index: int, total: int, layout_type: str = 'standard') -> tuple:
    """
    Calculate chart position based on index, total number of charts, and layout type.
    """
    if layout_type == 'with_highlights':
        # Chart on left side when highlights panel is present
        return Inches(0.5), Inches(1.5), Inches(5.5), Inches(4)
    
    if total == 1:
        # Center the chart
        return Inches(1), Inches(2), Inches(8), Inches(4)
    elif total == 2:
        # Side by side
        if index == 0:
            return Inches(0.5), Inches(2), Inches(4.5), Inches(4)
        else:
            return Inches(5), Inches(2), Inches(4.5), Inches(4)
    else:
        # Grid layout
        col = index % 2
        row = index // 2
        x = Inches(0.5 + col * 4.5)
        y = Inches(1.5 + row * 2.5)
        return x, y, Inches(4), Inches(2)

def add_bar_chart(slide, chart_spec: Dict[str, Any], x, y, cx, cy):
    """
    Add a bar chart to the slide.
    """
    try:
        chart_data = ChartData()
        chart_data.categories = chart_spec.get('data', {}).get('categories', [])
        
        series_data = chart_spec.get('data', {}).get('series', [])
        for series in series_data:
            chart_data.add_series(series.get('name', ''), series.get('values', []))
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            x, y, cx, cy,
            chart_data
        ).chart
        
        # Customize chart
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_spec.get('title', '')
        
    except Exception as e:
        logger.error(f"Error adding bar chart: {str(e)}")

def add_line_chart(slide, chart_spec: Dict[str, Any], x, y, cx, cy):
    """
    Add a line chart to the slide.
    """
    try:
        chart_data = ChartData()
        chart_data.categories = chart_spec.get('data', {}).get('categories', [])
        
        series_data = chart_spec.get('data', {}).get('series', [])
        for series in series_data:
            chart_data.add_series(series.get('name', ''), series.get('values', []))
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE,
            x, y, cx, cy,
            chart_data
        ).chart
        
        # Customize chart
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_spec.get('title', '')
        
    except Exception as e:
        logger.error(f"Error adding line chart: {str(e)}")

def add_pie_chart(slide, chart_spec: Dict[str, Any], x, y, cx, cy):
    """
    Add a pie chart to the slide.
    """
    try:
        chart_data = ChartData()
        chart_data.categories = chart_spec.get('data', {}).get('categories', [])
        chart_data.add_series('', chart_spec.get('data', {}).get('values', []))
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            x, y, cx, cy,
            chart_data
        ).chart
        
        # Customize chart
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_spec.get('title', '')
        
    except Exception as e:
        logger.error(f"Error adding pie chart: {str(e)}")

def add_donut_chart(slide, chart_spec: Dict[str, Any], x, y, cx, cy):
    """
    Add a donut chart to the slide.
    """
    try:
        chart_data = ChartData()
        chart_data.categories = chart_spec.get('data', {}).get('categories', [])
        chart_data.add_series('', chart_spec.get('data', {}).get('values', []))
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            x, y, cx, cy,
            chart_data
        ).chart
        
        # Customize chart
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_spec.get('title', '')
        
        # Add center text if specified
        center_text = chart_spec.get('center_text')
        if center_text:
            # Note: Adding text in the center of a donut chart requires additional shape manipulation
            # This is a placeholder for the actual implementation
            pass
        
    except Exception as e:
        logger.error(f"Error adding donut chart: {str(e)}")

def add_combo_chart(slide, chart_spec: Dict[str, Any], x, y, cx, cy):
    """
    Add a combination chart (bar + line) to the slide.
    """
    try:
        chart_data = ChartData()
        chart_data.categories = chart_spec.get('data', {}).get('categories', [])
        
        # Add bar series
        bar_series = chart_spec.get('data', {}).get('bar_series', [])
        for series in bar_series:
            chart_data.add_series(series.get('name', ''), series.get('values', []))
        
        # Create chart with bars
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            x, y, cx, cy,
            chart_data
        ).chart
        
        # Add line series on secondary axis
        line_series = chart_spec.get('data', {}).get('line_series', [])
        if line_series:
            # Create a secondary value axis
            chart.has_secondary_value_axis = True
            
            # Add line series
            for i, series in enumerate(line_series):
                line_data = ChartData()
                line_data.categories = chart_spec.get('data', {}).get('categories', [])
                line_data.add_series(series.get('name', ''), series.get('values', []))
                
                # Add the line chart overlay
                # Note: python-pptx has limitations with true combo charts
                # This is a simplified implementation
        
        # Customize chart
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_spec.get('title', '')
        
        # Style customization
        if chart_spec.get('style'):
            apply_chart_styling(chart, chart_spec.get('style'))
        
    except Exception as e:
        logger.error(f"Error adding combo chart: {str(e)}")

def process_table_slide(slide, content: Dict[str, Any]):
    """
    Process slides with tables.
    """
    try:
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = content.get('title', '')
        
        # Add tables
        tables_data = content.get('tables', [])
        add_tables_to_slide(slide, tables_data)
        
    except Exception as e:
        logger.error(f"Error processing table slide: {str(e)}")

def add_tables_to_slide(slide, tables_data: List[Dict[str, Any]]):
    """
    Add tables to slide.
    """
    try:
        for i, table_spec in enumerate(tables_data):
            headers = table_spec.get('headers', [])
            rows = table_spec.get('rows', [])
            
            if not headers or not rows:
                continue
            
            # Calculate position
            x = Inches(0.5)
            y = Inches(2 + i * 2.5)
            cx = Inches(9)
            cy = Inches(2)
            
            # Create table
            table = slide.shapes.add_table(
                rows=len(rows) + 1,
                cols=len(headers),
                left=x,
                top=y,
                width=cx,
                height=cy
            ).table
            
            # Set headers
            for j, header in enumerate(headers):
                cell = table.cell(0, j)
                cell.text = str(header)
                # Style header
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                paragraph = cell.text_frame.paragraphs[0]
                font = paragraph.font
                font.bold = True
                font.color.rgb = RGBColor(255, 255, 255)
            
            # Set data
            for i, row in enumerate(rows):
                for j, value in enumerate(row):
                    cell = table.cell(i + 1, j)
                    cell.text = str(value)
                    
    except Exception as e:
        logger.error(f"Error adding tables to slide: {str(e)}")

def process_content_slide(slide, content: Dict[str, Any]):
    """
    Process generic content slide.
    """
    try:
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = content.get('title', '')
        
        # Add content
        content_data = content.get('content', {})
        
        for shape in slide.placeholders:
            if hasattr(shape, 'text_frame') and shape.placeholder_format.idx > 0:
                tf = shape.text_frame
                tf.clear()
                
                # Add content based on structure
                if isinstance(content_data, str):
                    tf.text = content_data
                elif isinstance(content_data, list):
                    for item in content_data:
                        p = tf.add_paragraph()
                        p.text = f"• {item}"
                        p.level = 0
                elif isinstance(content_data, dict):
                    for key, value in content_data.items():
                        p = tf.add_paragraph()
                        p.text = f"{key}:"
                        p.level = 0
                        
                        if isinstance(value, list):
                            for item in value:
                                p = tf.add_paragraph()
                                p.text = f"• {item}"
                                p.level = 1
                        else:
                            p = tf.add_paragraph()
                            p.text = str(value)
                            p.level = 1
                
                break
                
    except Exception as e:
        logger.error(f"Error processing content slide: {str(e)}")

def process_generic_slide(slide, content: Dict[str, Any]):
    """
    Process any other slide type generically.
    """
    try:
        # Set title if available
        if slide.shapes.title and content.get('title'):
            slide.shapes.title.text = content['title']
        
        # Try to add content to first available text placeholder
        for shape in slide.placeholders:
            if hasattr(shape, 'text_frame') and shape.placeholder_format.idx > 0:
                content_str = json.dumps(content.get('content', ''), indent=2)
                shape.text = content_str
                break
                
    except Exception as e:
        logger.error(f"Error processing generic slide: {str(e)}")

def upload_presentation(buffer: io.BytesIO, output_key: str):
    """
    Upload presentation to S3.
    """
    try:
        s3.put_object(
            Bucket=OUTPUT_BUCKET,
            Key=output_key,
            Body=buffer.getvalue(),
            ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
        logger.info(f"Presentation uploaded to s3://{OUTPUT_BUCKET}/{output_key}")
        
    except Exception as e:
        logger.error(f"Error uploading presentation: {str(e)}")
        raise

def apply_corporate_styling(slide, style_config: Dict[str, Any]):
    """
    Apply corporate styling to slide including colors, logo, and footer.
    """
    try:
        # Corporate colors
        primary_color = style_config.get('primary_color', {'r': 255, 'g': 0, 'b': 0})  # Red default
        secondary_color = style_config.get('secondary_color', {'r': 128, 'g': 128, 'b': 128})  # Gray default
        
        # Add company logo to top right
        if style_config.get('logo_path'):
            try:
                logo_left = Inches(8.5)
                logo_top = Inches(0.3)
                logo_height = Inches(0.5)
                slide.shapes.add_picture(
                    style_config['logo_path'],
                    logo_left, logo_top,
                    height=logo_height
                )
            except Exception as e:
                logger.warning(f"Could not add logo: {str(e)}")
        
        # Add footer bar
        footer_height = Inches(0.5)
        footer_top = Inches(6.5)
        footer_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), footer_top,
            Inches(10), footer_height
        )
        footer_shape.fill.solid()
        footer_shape.fill.fore_color.rgb = RGBColor(
            secondary_color['r'], 
            secondary_color['g'], 
            secondary_color['b']
        )
        footer_shape.line.fill.background()
        
        # Add footer text
        footer_text = style_config.get('footer_text', 'South Plains Financial, Inc.')
        text_box = slide.shapes.add_textbox(
            Inches(0.5), footer_top + Inches(0.1),
            Inches(9), Inches(0.3)
        )
        text_frame = text_box.text_frame
        text_frame.text = footer_text
        text_frame.paragraphs[0].font.size = Pt(10)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Apply accent colors to charts
        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart
                # Apply primary color to first series
                if chart.series:
                    chart.series[0].format.fill.solid()
                    chart.series[0].format.fill.fore_color.rgb = RGBColor(
                        primary_color['r'],
                        primary_color['g'], 
                        primary_color['b']
                    )
                    
    except Exception as e:
        logger.error(f"Error applying corporate styling: {str(e)}")

def apply_chart_styling(chart, style_config: Dict[str, Any]):
    """
    Apply styling to charts including colors and formatting.
    """
    try:
        # Apply colors to series
        colors = style_config.get('series_colors', [])
        for i, series in enumerate(chart.series):
            if i < len(colors):
                color = colors[i]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = RGBColor(
                    color.get('r', 0),
                    color.get('g', 0),
                    color.get('b', 0)
                )
        
        # Apply line styles for line series
        line_styles = style_config.get('line_styles', [])
        for i, series in enumerate(chart.series):
            if hasattr(series, 'format') and hasattr(series.format, 'line'):
                if i < len(line_styles):
                    style = line_styles[i]
                    if style.get('dashed'):
                        series.format.line.dash_style = 'sysDash'
                        
    except Exception as e:
        logger.error(f"Error applying chart styling: {str(e)}")

def create_slide_with_highlights_layout(slide, content: Dict[str, Any]):
    """
    Create a slide with chart on left and highlights panel on right.
    """
    try:
        # Chart area (left side)
        chart_left = Inches(0.5)
        chart_top = Inches(1.5)
        chart_width = Inches(6)
        chart_height = Inches(4)
        
        # Highlights panel (right side)
        highlights_left = Inches(6.8)
        highlights_top = Inches(1.5)
        highlights_width = Inches(3)
        highlights_height = Inches(4)
        
        # Add highlights panel background
        panel_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            highlights_left, highlights_top,
            highlights_width, highlights_height
        )
        panel_shape.fill.solid()
        panel_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray
        panel_shape.line.fill.background()
        
        # Add highlights title
        title_box = slide.shapes.add_textbox(
            highlights_left + Inches(0.1), 
            highlights_top + Inches(0.1),
            highlights_width - Inches(0.2), 
            Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get('highlights_title', '2Q\'20 Highlights')
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.size = Pt(14)
        
        # Add highlights content
        content_box = slide.shapes.add_textbox(
            highlights_left + Inches(0.1),
            highlights_top + Inches(0.7),
            highlights_width - Inches(0.2),
            highlights_height - Inches(0.8)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        highlights = content.get('highlights', [])
        for highlight in highlights:
            p = content_frame.add_paragraph()
            p.text = f"• {highlight}"
            p.font.size = Pt(11)
            p.space_after = Pt(6)
            
        return chart_left, chart_top, chart_width, chart_height
        
    except Exception as e:
        logger.error(f"Error creating slide with highlights layout: {str(e)}")
        return None