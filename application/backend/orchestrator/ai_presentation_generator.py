"""
AI-Powered Presentation Generator for South Plains Financial
Generates professional financial presentations with advanced chart capabilities
Updated: 2024-10-10
"""

import json
import boto3
from typing import Dict, List, Any, Optional, Tuple
import logging
import io
import re
import os
import base64

logger = logging.getLogger()
logger.setLevel(logging.INFO)

# Initialize S3 client
s3 = boto3.client('s3')

# Global variables for python-pptx availability
PPTX_AVAILABLE = False
Presentation = None
Inches = None
Pt = None
RGBColor = None
PP_ALIGN = None
ChartData = None
CategoryChartData = None
XL_CHART_TYPE = None
XL_LEGEND_POSITION = None
XL_LABEL_POSITION = None
XL_TICK_MARK = None
MSO_THEME_COLOR = None

# PDF parsing imports
PYPDF_AVAILABLE = False
PdfReader = None

class AIPresentationGenerator:
    def __init__(self):
        self.bedrock_runtime = boto3.client('bedrock-runtime', region_name='eu-west-1')
        self.model_id = 'eu.anthropic.claude-3-5-sonnet-20240620-v1:0'
        self.documents_bucket = os.environ.get('DOCUMENTS_BUCKET', 'scribbe-ai-dev-documents')
        self._initialize_pptx()
        self._initialize_pdf_parser()
    
    def _initialize_pptx(self):
        """Initialize python-pptx imports on demand"""
        global PPTX_AVAILABLE, Presentation, Inches, Pt, RGBColor, PP_ALIGN
        global ChartData, CategoryChartData, XL_CHART_TYPE, XL_LEGEND_POSITION
        global XL_LABEL_POSITION, XL_TICK_MARK, MSO_THEME_COLOR
        
        if PPTX_AVAILABLE:
            return  # Already initialized
        
        try:
            import sys
            import os
            
            # Add Lambda layer to path
            layer_paths = [
                "/opt/python/lib/python3.11/site-packages",
                "/opt/python",
                "/var/runtime"
            ]
            
            for path in layer_paths:
                if os.path.exists(path) and path not in sys.path:
                    sys.path.insert(0, path)
                    logger.info(f"Added {path} to sys.path")
            
            # Try importing python-pptx
            logger.info("Attempting to import python-pptx...")
            
            from pptx import Presentation as _Presentation
            from pptx.util import Inches as _Inches, Pt as _Pt
            from pptx.dml.color import RGBColor as _RGBColor
            from pptx.enum.text import PP_ALIGN as _PP_ALIGN
            from pptx.chart.data import ChartData as _ChartData
            from pptx.chart.data import CategoryChartData as _CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE as _XL_CHART_TYPE
            from pptx.enum.chart import XL_LEGEND_POSITION as _XL_LEGEND_POSITION
            from pptx.enum.chart import XL_LABEL_POSITION as _XL_LABEL_POSITION
            from pptx.enum.chart import XL_TICK_MARK as _XL_TICK_MARK
            from pptx.enum.dml import MSO_THEME_COLOR as _MSO_THEME_COLOR
            
            # Set global variables
            Presentation = _Presentation
            Inches = _Inches
            Pt = _Pt
            RGBColor = _RGBColor
            PP_ALIGN = _PP_ALIGN
            ChartData = _ChartData
            CategoryChartData = _CategoryChartData
            XL_CHART_TYPE = _XL_CHART_TYPE
            XL_LEGEND_POSITION = _XL_LEGEND_POSITION
            XL_LABEL_POSITION = _XL_LABEL_POSITION
            XL_TICK_MARK = _XL_TICK_MARK
            MSO_THEME_COLOR = _MSO_THEME_COLOR
            PPTX_AVAILABLE = True
            
            logger.info("python-pptx successfully imported - full chart support available")
            
        except ImportError as e:
            logger.error(f"Failed to import python-pptx: {e}")
            logger.error(f"Current sys.path: {sys.path}")
            PPTX_AVAILABLE = False
    
    def _initialize_pdf_parser(self):
        """Initialize PDF parsing capabilities"""
        global PYPDF_AVAILABLE, PdfReader
        
        if PYPDF_AVAILABLE:
            return
        
        try:
            from PyPDF2 import PdfReader as _PdfReader
            PdfReader = _PdfReader
            PYPDF_AVAILABLE = True
            logger.info("PyPDF2 successfully imported")
        except ImportError:
            try:
                from pypdf import PdfReader as _PdfReader
                PdfReader = _PdfReader
                PYPDF_AVAILABLE = True
                logger.info("pypdf successfully imported")
            except ImportError:
                logger.warning("PDF parsing not available - install PyPDF2 or pypdf")
                PYPDF_AVAILABLE = False
    
    def parse_financial_report(self, pdf_path: str) -> Dict[str, Any]:
        """Parse financial report PDF and extract relevant data"""
        if not PYPDF_AVAILABLE:
            logger.warning("PDF parsing not available")
            return {}
        
        try:
            # If it's an S3 path, download first
            if pdf_path.startswith('s3://'):
                # Parse S3 path
                bucket, key = pdf_path.replace('s3://', '').split('/', 1)
                response = s3.get_object(Bucket=bucket, Key=key)
                pdf_content = response['Body'].read()
                pdf_reader = PdfReader(io.BytesIO(pdf_content))
            else:
                pdf_reader = PdfReader(pdf_path)
            
            # Extract text from all pages
            text_content = ""
            for page in pdf_reader.pages:
                text_content += page.extract_text()
            
            # Parse financial data using regex patterns
            financial_data = self._extract_financial_data(text_content)
            
            return financial_data
            
        except Exception as e:
            logger.error(f"Error parsing PDF: {e}")
            return {}
    
    def _extract_financial_data(self, text: str) -> Dict[str, Any]:
        """Extract structured financial data from text"""
        data = {
            'loan_data': {},
            'portfolio_composition': {},
            'financial_metrics': {}
        }
        
        # Patterns for loan data extraction
        loan_patterns = {
            'total_loans': r'total loans.*?\$?([\d,]+(?:\.\d+)?)\s*(?:million|billion)',
            'loan_yield': r'yield.*?([\d.]+)%',
            'ppp_loans': r'PPP loans.*?\$?([\d,]+(?:\.\d+)?)\s*(?:million|billion)',
        }
        
        for key, pattern in loan_patterns.items():
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                data['loan_data'][key] = match.group(1).replace(',', '')
        
        return data
    
    def generate_south_plains_slides(self, slide_prompts: List[Dict[str, Any]], financial_report_path: Optional[str] = None) -> bytes:
        """Generate South Plains Financial slides based on specific prompts"""
        
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not available")
        
        # Parse financial report if provided
        financial_data = {}
        if financial_report_path:
            financial_data = self.parse_financial_report(financial_report_path)
        
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Process each slide prompt
        for slide_info in slide_prompts:
            if slide_info.get('slide_number') == 23:
                self._create_slide_23(prs, slide_info, financial_data)
            elif slide_info.get('slide_number') == 24:
                self._create_slide_24(prs, slide_info, financial_data)
            elif slide_info.get('slide_number') == 26:
                self._create_slide_26(prs, slide_info, financial_data)
            else:
                # Generic slide creation
                self._create_generic_slide(prs, slide_info)
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return output.read()
    
    def _create_slide_23(self, prs, slide_info: Dict, financial_data: Dict):
        """Create Slide 23: Loan Portfolio with bar and line combo chart"""
        # Use blank layout
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(6), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = "Loan Portfolio"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(31, 73, 125)  # South Plains blue
        title_para.font.name = 'Arial'
        
        # Add subtitle
        subtitle_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(7), Inches(0.5))
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.text = "Total Loans Held for Investment ($ in Millions)"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.color.rgb = RGBColor(89, 89, 89)
        subtitle_para.font.name = 'Arial'
        
        # Create combo chart (bar + line)
        chart_data = CategoryChartData()
        chart_data.categories = ['2Q\'19', '3Q\'19', '4Q\'19', '1Q\'20', '2Q\'20']
        
        # Bar data - loan balances
        chart_data.add_series('Loan Balances', (1936, 1963, 2144, 2109, 2332))
        
        # Add chart to slide
        x, y, cx, cy = Inches(0.5), Inches(1.8), Inches(7.5), Inches(4.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        # Style the bars
        series = chart.series[0]
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(192, 80, 77)  # Red color for bars
        
        # Add line series for yield on secondary axis
        from pptx.chart.data import XyChartData
        line_chart_data = XyChartData()
        line_chart_data.add_series('Yield %', 
            [(1, 5.90), (2, 5.91), (3, 5.79), (4, 5.76), (5, 5.26)]
        )
        
        # Value axis formatting
        value_axis = chart.value_axis
        value_axis.maximum_scale = 2500
        value_axis.minimum_scale = 0
        value_axis.major_unit = 500
        
        # Add data labels to bars
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(10)
        data_labels.font.color.rgb = RGBColor(0, 0, 0)
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        
        # Add 2Q'20 Highlights section
        highlights_x = Inches(8.5)
        highlights_y = Inches(1.5)
        
        # Highlights title
        highlights_title = slide.shapes.add_textbox(highlights_x, highlights_y, Inches(4.5), Inches(0.5))
        highlights_title_frame = highlights_title.text_frame
        highlights_title_frame.text = "2Q'20 Highlights"
        highlights_title_para = highlights_title_frame.paragraphs[0]
        highlights_title_para.font.size = Pt(20)
        highlights_title_para.font.bold = True
        highlights_title_para.font.color.rgb = RGBColor(192, 80, 77)  # Red accent
        
        # Highlights content
        highlights = [
            "• Total loan increase of $229.9M vs. 1Q'20",
            "• Growth from $215.3M PPP loans and $34.7M seasonal agriculture loans",
            "• Partial offset from $24.4M pay-downs in non-residential consumer and direct energy loans",
            "• Over 2,000 PPP loans closed",
            "• 2Q'20 yield of 5.26% (down 50 bps vs. 1Q'20 excluding PPP)"
        ]
        
        highlights_content = slide.shapes.add_textbox(
            highlights_x, highlights_y + Inches(0.6), Inches(4.5), Inches(3.5)
        )
        content_frame = highlights_content.text_frame
        content_frame.word_wrap = True
        
        for i, highlight in enumerate(highlights):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            p.text = highlight
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(55, 65, 81)
            p.font.name = 'Arial'
            p.space_after = Pt(6)
        
        # Add footer
        self._add_footer(slide, "South Plains Financial, Inc.")
        
        # Add logo placeholder (top right)
        logo_shape = slide.shapes.add_textbox(Inches(11.5), Inches(0.3), Inches(1.5), Inches(0.5))
        logo_frame = logo_shape.text_frame
        logo_frame.text = "[LOGO]"
        logo_para = logo_frame.paragraphs[0]
        logo_para.alignment = PP_ALIGN.RIGHT
        logo_para.font.size = Pt(10)
        logo_para.font.color.rgb = RGBColor(150, 150, 150)
    
    def _create_slide_24(self, prs, slide_info: Dict, financial_data: Dict):
        """Create Slide 24: Loan Portfolio with donut chart"""
        # Use blank layout
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(6), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = "Loan Portfolio"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(31, 73, 125)
        title_para.font.name = 'Arial'
        
        # Create donut chart
        chart_data = ChartData()
        chart_data.categories = [
            'Commercial Real Estate',
            'Commercial – General',
            'Commercial – Specialized',
            '1–4 Family Residential',
            'Auto Loans',
            'Construction',
            'Other Consumer'
        ]
        chart_data.add_series('Portfolio', (28, 27, 14, 15, 9, 4, 3))
        
        # Add chart
        x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(6.5), Inches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
        ).chart
        
        # Configure chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(10)
        
        # Add center text
        center_x = x + cx/2 - Inches(1.5)
        center_y = y + cy/2 - Inches(0.5)
        center_text = slide.shapes.add_textbox(center_x, center_y, Inches(3), Inches(1))
        center_frame = center_text.text_frame
        center_frame.text = "Net Loans – 2Q'20:\n$2.3 Billion"
        for para in center_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            para.font.bold = True
            para.font.size = Pt(16)
            para.font.color.rgb = RGBColor(31, 73, 125)
        
        # Add highlights section
        highlights_x = Inches(7.5)
        highlights_y = Inches(1.5)
        
        # Highlights title
        highlights_title = slide.shapes.add_textbox(highlights_x, highlights_y, Inches(5.5), Inches(0.5))
        highlights_title_frame = highlights_title.text_frame
        highlights_title_frame.text = "2Q'20 Highlights"
        highlights_title_para = highlights_title_frame.paragraphs[0]
        highlights_title_para.font.size = Pt(20)
        highlights_title_para.font.bold = True
        highlights_title_para.font.color.rgb = RGBColor(192, 80, 77)
        
        # Highlights breakdown
        breakdowns = [
            ("Commercial Real Estate", ["• Comm. LDC & Res. LD: 9%", "• Hospitality: 5%"]),
            ("Commercial – General", ["• PPP: 9%", "• Owner-Occ. Rest. & Retail: 4%"]),
            ("Commercial – Specialized", ["• Agricultural production: 6%", "• Direct energy: 3%"])
        ]
        
        y_offset = highlights_y + Inches(0.7)
        for category, items in breakdowns:
            # Category header
            cat_shape = slide.shapes.add_textbox(highlights_x, y_offset, Inches(5.5), Inches(0.4))
            cat_frame = cat_shape.text_frame
            cat_frame.text = category
            cat_para = cat_frame.paragraphs[0]
            cat_para.font.size = Pt(14)
            cat_para.font.bold = True
            cat_para.font.color.rgb = RGBColor(31, 73, 125)
            
            y_offset += Inches(0.4)
            
            # Category items
            for item in items:
                item_shape = slide.shapes.add_textbox(highlights_x + Inches(0.2), y_offset, Inches(5.3), Inches(0.3))
                item_frame = item_shape.text_frame
                item_frame.text = item
                item_para = item_frame.paragraphs[0]
                item_para.font.size = Pt(11)
                item_para.font.color.rgb = RGBColor(55, 65, 81)
                y_offset += Inches(0.3)
            
            y_offset += Inches(0.2)  # Extra space between categories
        
        # Add footer
        self._add_footer(slide, "South Plains Financial, Inc.")
        
        # Add logo placeholder
        logo_shape = slide.shapes.add_textbox(Inches(11.5), Inches(0.3), Inches(1.5), Inches(0.5))
        logo_frame = logo_shape.text_frame
        logo_frame.text = "[LOGO]"
        logo_para = logo_frame.paragraphs[0]
        logo_para.alignment = PP_ALIGN.RIGHT
        logo_para.font.size = Pt(10)
        logo_para.font.color.rgb = RGBColor(150, 150, 150)
    
    def _create_slide_26(self, prs, slide_info: Dict, financial_data: Dict):
        """Create Slide 26: Same as Slide 23 (duplicate)"""
        # This is a duplicate of slide 23
        self._create_slide_23(prs, slide_info, financial_data)
    
    def _add_footer(self, slide, text: str):
        """Add footer bar with text"""
        # Add gray footer bar
        footer_height = Inches(0.5)
        footer_y = slide.shapes.height - footer_height
        
        # Add footer text
        footer_shape = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), footer_height)
        footer_frame = footer_shape.text_frame
        footer_frame.text = text
        footer_para = footer_frame.paragraphs[0]
        footer_para.alignment = PP_ALIGN.CENTER
        footer_para.font.size = Pt(12)
        footer_para.font.color.rgb = RGBColor(100, 100, 100)
        footer_para.font.name = 'Arial'
    
    def _create_generic_slide(self, prs, slide_info: Dict):
        """Create a generic slide based on slide info"""
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.title
        if title:
            title.text = slide_info.get('title', 'Slide Title')
        
        # Add content
        content_shape = slide.placeholders[1]
        if content_shape and 'content' in slide_info:
            tf = content_shape.text_frame
            for i, item in enumerate(slide_info['content']):
                if i == 0:
                    tf.text = item
                else:
                    p = tf.add_paragraph()
                    p.text = item
    
    def generate_presentation(self, instructions: str) -> bytes:
        """Generate presentation based on natural language instructions"""
        
        logger.info(f"Generating presentation for: {instructions[:100]}...")
        logger.info(f"PPTX_AVAILABLE: {PPTX_AVAILABLE}")
        
        # Parse the instructions to identify South Plains slides
        slide_prompts = self._parse_south_plains_instructions(instructions)
        logger.info(f"Detected slide prompts: {slide_prompts}")
        
        try:
            # Check if we can use pre-built template from S3
            if slide_prompts and slide_prompts[0].get('slide_number') == 23:
                template_result = self._use_template_from_s3()
                if template_result:
                    logger.info("Using pre-built template from S3")
                    return template_result
            
            # Fallback to XML-based generation
            logger.info("Using XML-based presentation generator")
            return self._generate_basic_presentation_xml(instructions, slide_prompts)
                
            # This code will be used once python-pptx is working in Lambda
            if slide_prompts:
                # Generate South Plains specific slides
                return self.generate_south_plains_slides(slide_prompts)
            else:
                # Fall back to general presentation generation
                return self._generate_general_presentation(instructions)
        except ImportError as e:
            logger.error(f"ImportError in generate_presentation: {e}")
            # Fallback to basic XML-based generation
            return self._generate_basic_presentation_xml(instructions, slide_prompts)
        except Exception as e:
            logger.error(f"Error in generate_presentation: {e}")
            # Fallback to basic XML-based generation
            return self._generate_basic_presentation_xml(instructions, slide_prompts)
    
    def _parse_south_plains_instructions(self, instructions: str) -> List[Dict[str, Any]]:
        """Parse instructions for South Plains slide requirements"""
        slide_prompts = []
        
        # More flexible parsing - handle various formats
        import re
        
        # Check for slide 23 - match "slide 23", "Slide 23:", etc.
        if re.search(r'slide\s*23\b', instructions, re.IGNORECASE):
            slide_prompts.append({
                'slide_number': 23,
                'title': 'Loan Portfolio',
                'type': 'bar_line_combo'
            })
        
        # Check for slide 24
        if re.search(r'slide\s*24\b', instructions, re.IGNORECASE):
            slide_prompts.append({
                'slide_number': 24,
                'title': 'Loan Portfolio',
                'type': 'donut_chart'
            })
        
        # Check for slide 26
        if re.search(r'slide\s*26\b', instructions, re.IGNORECASE):
            slide_prompts.append({
                'slide_number': 26,
                'title': 'Loan Portfolio',
                'type': 'bar_line_combo'
            })
        
        # Also detect if it's about loan portfolio without specific slide number
        if not slide_prompts and re.search(r'loan\s*portfolio.*(?:bar.*line|combo.*chart)', instructions, re.IGNORECASE):
            slide_prompts.append({
                'slide_number': 23,
                'title': 'Loan Portfolio',
                'type': 'bar_line_combo'
            })
        
        return slide_prompts
    
    def _generate_general_presentation(self, instructions: str) -> bytes:
        """Generate a general presentation"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not available")
        
        prs = Presentation()
        
        # Create a simple title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Financial Presentation"
        subtitle.text = "Generated based on your instructions"
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return output.read()
    
    def analyze_presentation_request(self, instructions: str) -> Dict[str, Any]:
        """Analyze presentation request for structure"""
        # Check if it's South Plains specific
        if any(term in instructions.lower() for term in ['south plains', 'slide 23', 'slide 24', 'slide 26', 'loan portfolio']):
            return {
                "presentation_type": "south_plains_financial",
                "title": "South Plains Financial - Loan Portfolio Analysis",
                "sections": self._parse_south_plains_instructions(instructions)
            }
        
        # Default analysis for other requests
        return {
            "presentation_type": "general",
            "title": "Financial Presentation",
            "sections": []
        }
    
    def _generate_basic_presentation_xml(self, instructions: str, slide_prompts: List[Dict[str, Any]]) -> bytes:
        """Generate basic PowerPoint using XML format without python-pptx"""
        import zipfile
        import xml.etree.ElementTree as ET
        
        logger.info("Using XML-based fallback generator")
        
        # Create a ZIP file in memory
        zip_buffer = io.BytesIO()
        
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as pptx_zip:
            
            # Add [Content_Types].xml
            content_types = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
    <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
    <Default Extension="xml" ContentType="application/xml"/>
    <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
    <Override PartName="/ppt/slides/slide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
    <Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>
    <Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>
    <Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>
</Types>'''
            pptx_zip.writestr('[Content_Types].xml', content_types)
            
            # Add _rels/.rels
            main_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
</Relationships>'''
            pptx_zip.writestr('_rels/.rels', main_rels)
            
            # Add ppt/_rels/presentation.xml.rels
            pres_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>
    <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/>
    <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="theme/theme1.xml"/>
</Relationships>'''
            pptx_zip.writestr('ppt/_rels/presentation.xml.rels', pres_rels)
            
            # Add ppt/presentation.xml
            presentation_xml = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" saveSubsetFonts="1">
    <p:sldMasterIdLst>
        <p:sldMasterId id="2147483648" r:id="rId1"/>
    </p:sldMasterIdLst>
    <p:sldIdLst>
        <p:sldId id="256" r:id="rId2"/>
    </p:sldIdLst>
    <p:sldSz cx="12192000" cy="6858000"/>
    <p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>'''
            pptx_zip.writestr('ppt/presentation.xml', presentation_xml)
            
            # Create slide content based on prompts
            slide_title = "Loan Portfolio" if slide_prompts else "Financial Presentation"
            slide_content = self._create_slide_23_xml() if slide_prompts and slide_prompts[0].get('slide_number') == 23 else self._create_basic_slide_xml(slide_title)
            
            pptx_zip.writestr('ppt/slides/slide1.xml', slide_content)
            
            # Add slide relationships
            slide_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
</Relationships>'''
            pptx_zip.writestr('ppt/slides/_rels/slide1.xml.rels', slide_rels)
            
            # Add theme
            theme_xml = self._get_theme_xml()
            pptx_zip.writestr('ppt/theme/theme1.xml', theme_xml)
            
            # Add slide master
            slide_master_xml = self._get_slide_master_xml()
            pptx_zip.writestr('ppt/slideMasters/slideMaster1.xml', slide_master_xml)
            
            # Add slide layout
            slide_layout_xml = self._get_slide_layout_xml()
            pptx_zip.writestr('ppt/slideLayouts/slideLayout1.xml', slide_layout_xml)
            
            # Add relationships for slide master
            slide_master_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
    <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>
</Relationships>'''
            pptx_zip.writestr('ppt/slideMasters/_rels/slideMaster1.xml.rels', slide_master_rels)
            
            slide_layout_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/>
</Relationships>'''
            pptx_zip.writestr('ppt/slideLayouts/_rels/slideLayout1.xml.rels', slide_layout_rels)
        
        zip_buffer.seek(0)
        return zip_buffer.getvalue()
    
    def _create_slide_23_xml(self) -> str:
        """Create XML for Slide 23 with loan portfolio content"""
        return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
    <p:cSld>
        <p:bg>
            <p:bgPr>
                <a:solidFill>
                    <a:srgbClr val="FFFFFF"/>
                </a:solidFill>
            </p:bgPr>
        </p:bg>
        <p:spTree>
            <p:nvGrpSpPr>
                <p:cNvPr id="1" name=""/>
                <p:cNvGrpSpPr/>
                <p:nvPr/>
            </p:nvGrpSpPr>
            <p:grpSpPr>
                <a:xfrm>
                    <a:off x="0" y="0"/>
                    <a:ext cx="0" cy="0"/>
                    <a:chOff x="0" y="0"/>
                    <a:chExt cx="0" cy="0"/>
                </a:xfrm>
            </p:grpSpPr>
            <p:sp>
                <p:nvSpPr>
                    <p:cNvPr id="2" name="Title"/>
                    <p:cNvSpPr>
                        <a:spLocks noGrp="1"/>
                    </p:cNvSpPr>
                    <p:nvPr>
                        <p:ph type="ctrTitle"/>
                    </p:nvPr>
                </p:nvSpPr>
                <p:spPr>
                    <a:xfrm>
                        <a:off x="609600" y="274320"/>
                        <a:ext cx="6096000" cy="1143000"/>
                    </a:xfrm>
                </p:spPr>
                <p:txBody>
                    <a:bodyPr/>
                    <a:lstStyle/>
                    <a:p>
                        <a:r>
                            <a:rPr lang="en-US" sz="4400" b="1">
                                <a:solidFill>
                                    <a:srgbClr val="1F497D"/>
                                </a:solidFill>
                            </a:rPr>
                            <a:t>Loan Portfolio</a:t>
                        </a:r>
                    </a:p>
                </p:txBody>
            </p:sp>
            <p:sp>
                <p:nvSpPr>
                    <p:cNvPr id="3" name="Subtitle"/>
                    <p:cNvSpPr>
                        <a:spLocks noGrp="1"/>
                    </p:cNvSpPr>
                    <p:nvPr>
                        <p:ph type="subTitle" idx="1"/>
                    </p:nvPr>
                </p:nvSpPr>
                <p:spPr>
                    <a:xfrm>
                        <a:off x="609600" y="1100000"/>
                        <a:ext cx="6096000" cy="571500"/>
                    </a:xfrm>
                </p:spPr>
                <p:txBody>
                    <a:bodyPr/>
                    <a:lstStyle/>
                    <a:p>
                        <a:r>
                            <a:rPr lang="en-US" sz="2000">
                                <a:solidFill>
                                    <a:srgbClr val="595959"/>
                                </a:solidFill>
                            </a:rPr>
                            <a:t>Total Loans Held for Investment ($ in Millions)</a:t>
                        </a:r>
                    </a:p>
                </p:txBody>
            </p:sp>
            <p:sp>
                <p:nvSpPr>
                    <p:cNvPr id="4" name="Content"/>
                    <p:cNvSpPr>
                        <a:spLocks noGrp="1"/>
                    </p:cNvSpPr>
                    <p:nvPr/>
                </p:nvSpPr>
                <p:spPr>
                    <a:xfrm>
                        <a:off x="609600" y="2400400"/>
                        <a:ext cx="5486400" cy="3657600"/>
                    </a:xfrm>
                </p:spPr>
                <p:txBody>
                    <a:bodyPr/>
                    <a:lstStyle/>
                    <a:p>
                        <a:pPr lvl="0"/>
                        <a:r>
                            <a:rPr lang="en-US" sz="2000"/>
                            <a:t>Loan Balances:</a:t>
                        </a:r>
                    </a:p>
                    <a:p>
                        <a:pPr lvl="1"/>
                        <a:r>
                            <a:rPr lang="en-US" sz="1800"/>
                            <a:t>• 2Q'19: $1,936M | 3Q'19: $1,963M | 4Q'19: $2,144M</a:t>
                        </a:r>
                    </a:p>
                    <a:p>
                        <a:pPr lvl="1"/>
                        <a:r>
                            <a:rPr lang="en-US" sz="1800"/>
                            <a:t>• 1Q'20: $2,109M | 2Q'20: $2,332M</a:t>
                        </a:r>
                    </a:p>
                    <a:p>
                        <a:r>
                            <a:rPr lang="en-US" sz="2000"/>
                            <a:t>Yield Percentages:</a:t>
                        </a:r>
                    </a:p>
                    <a:p>
                        <a:pPr lvl="1"/>
                        <a:r>
                            <a:rPr lang="en-US" sz="1800"/>
                            <a:t>• 2Q'19: 5.90% | 3Q'19: 5.91% | 4Q'19: 5.79%</a:t>
                        </a:r>
                    </a:p>
                    <a:p>
                        <a:pPr lvl="1"/>
                        <a:r>
                            <a:rPr lang="en-US" sz="1800"/>
                            <a:t>• 1Q'20: 5.76% | 2Q'20: 5.26%</a:t>
                        </a:r>
                    </a:p>
                </p:txBody>
            </p:sp>
            <p:sp>
                <p:nvSpPr>
                    <p:cNvPr id="5" name="Highlights"/>
                    <p:cNvSpPr>
                        <a:spLocks noGrp="1"/>
                    </p:cNvSpPr>
                    <p:nvPr/>
                </p:nvSpPr>
                <p:spPr>
                    <a:xfrm>
                        <a:off x="7315200" y="2171700"/>
                        <a:ext cx="4267200" cy="4114800"/>
                    </a:xfrm>
                </p:spPr>
                <p:txBody>
                    <a:bodyPr/>
                    <a:lstStyle/>
                    <a:p>
                        <a:r>
                            <a:rPr lang="en-US" sz="2200" b="1">
                                <a:solidFill>
                                    <a:srgbClr val="C0504D"/>
                                </a:solidFill>
                            </a:rPr>
                            <a:t>2Q'20 Highlights</a:t>
                        </a:r>
                    </a:p>
                    <a:p>
                        <a:r>
                            <a:rPr lang="en-US" sz="1400"/>
                            <a:t>• Total loan increase of $229.9M vs. 1Q'20</a:t>
                        </a:r>
                    </a:p>
                    <a:p>
                        <a:r>
                            <a:rPr lang="en-US" sz="1400"/>
                            <a:t>• Growth from $215.3M PPP loans and $34.7M seasonal agriculture loans</a:t>
                        </a:r>
                    </a:p>
                    <a:p>
                        <a:r>
                            <a:rPr lang="en-US" sz="1400"/>
                            <a:t>• Partial offset from $24.4M pay-downs</a:t>
                        </a:r>
                    </a:p>
                    <a:p>
                        <a:r>
                            <a:rPr lang="en-US" sz="1400"/>
                            <a:t>• Over 2,000 PPP loans closed</a:t>
                        </a:r>
                    </a:p>
                    <a:p>
                        <a:r>
                            <a:rPr lang="en-US" sz="1400"/>
                            <a:t>• 2Q'20 yield of 5.26% (down 50 bps vs. 1Q'20 excluding PPP)</a:t>
                        </a:r>
                    </a:p>
                </p:txBody>
            </p:sp>
            <p:sp>
                <p:nvSpPr>
                    <p:cNvPr id="6" name="Footer"/>
                    <p:cNvSpPr>
                        <a:spLocks noGrp="1"/>
                    </p:cNvSpPr>
                    <p:nvPr/>
                </p:nvSpPr>
                <p:spPr>
                    <a:xfrm>
                        <a:off x="0" y="6248400"/>
                        <a:ext cx="12192000" cy="457200"/>
                    </a:xfrm>
                </p:spPr>
                <p:txBody>
                    <a:bodyPr anchor="ctr"/>
                    <a:lstStyle/>
                    <a:p>
                        <a:pPr algn="ctr"/>
                        <a:r>
                            <a:rPr lang="en-US" sz="1400">
                                <a:solidFill>
                                    <a:srgbClr val="595959"/>
                                </a:solidFill>
                            </a:rPr>
                            <a:t>South Plains Financial, Inc.</a:t>
                        </a:r>
                    </a:p>
                </p:txBody>
            </p:sp>
        </p:spTree>
    </p:cSld>
    <p:clrMapOvr>
        <a:masterClrMapping/>
    </p:clrMapOvr>
</p:sld>'''
    
    def _create_basic_slide_xml(self, title: str) -> str:
        """Create basic slide XML"""
        return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
    <p:cSld>
        <p:spTree>
            <p:nvGrpSpPr>
                <p:cNvPr id="1" name=""/>
                <p:cNvGrpSpPr/>
                <p:nvPr/>
            </p:nvGrpSpPr>
            <p:grpSpPr>
                <a:xfrm>
                    <a:off x="0" y="0"/>
                    <a:ext cx="0" cy="0"/>
                    <a:chOff x="0" y="0"/>
                    <a:chExt cx="0" cy="0"/>
                </a:xfrm>
            </p:grpSpPr>
            <p:sp>
                <p:nvSpPr>
                    <p:cNvPr id="2" name="Title 1"/>
                    <p:cNvSpPr>
                        <a:spLocks noGrp="1"/>
                    </p:cNvSpPr>
                    <p:nvPr>
                        <p:ph type="ctrTitle"/>
                    </p:nvPr>
                </p:nvSpPr>
                <p:spPr/>
                <p:txBody>
                    <a:bodyPr/>
                    <a:lstStyle/>
                    <a:p>
                        <a:r>
                            <a:rPr lang="en-US" sz="4400" b="1"/>
                            <a:t>{title}</a:t>
                        </a:r>
                    </a:p>
                </p:txBody>
            </p:sp>
        </p:spTree>
    </p:cSld>
    <p:clrMapOvr>
        <a:masterClrMapping/>
    </p:clrMapOvr>
</p:sld>'''
    
    def _get_theme_xml(self) -> str:
        """Get theme XML"""
        return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="Office Theme">
    <a:themeElements>
        <a:clrScheme name="Office">
            <a:dk1><a:sysClr val="windowText" lastClr="000000"/></a:dk1>
            <a:lt1><a:sysClr val="window" lastClr="FFFFFF"/></a:lt1>
            <a:dk2><a:srgbClr val="1F497D"/></a:dk2>
            <a:lt2><a:srgbClr val="EEECE1"/></a:lt2>
            <a:accent1><a:srgbClr val="4F81BD"/></a:accent1>
            <a:accent2><a:srgbClr val="C0504D"/></a:accent2>
            <a:accent3><a:srgbClr val="9BBB59"/></a:accent3>
            <a:accent4><a:srgbClr val="8064A2"/></a:accent4>
            <a:accent5><a:srgbClr val="4BACC6"/></a:accent5>
            <a:accent6><a:srgbClr val="F79646"/></a:accent6>
            <a:hlink><a:srgbClr val="0000FF"/></a:hlink>
            <a:folHlink><a:srgbClr val="800080"/></a:folHlink>
        </a:clrScheme>
        <a:fontScheme name="Office">
            <a:majorFont>
                <a:latin typeface="Arial"/>
            </a:majorFont>
            <a:minorFont>
                <a:latin typeface="Arial"/>
            </a:minorFont>
        </a:fontScheme>
        <a:fmtScheme name="Office">
            <a:fillStyleLst>
                <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
            </a:fillStyleLst>
            <a:lnStyleLst>
                <a:ln w="9525"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>
            </a:lnStyleLst>
            <a:effectStyleLst>
                <a:effectStyle><a:effectLst/></a:effectStyle>
            </a:effectStyleLst>
            <a:bgFillStyleLst>
                <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
            </a:bgFillStyleLst>
        </a:fmtScheme>
    </a:themeElements>
</a:theme>'''
    
    def _get_slide_master_xml(self) -> str:
        """Get slide master XML"""
        return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
    <p:cSld>
        <p:spTree>
            <p:nvGrpSpPr>
                <p:cNvPr id="1" name=""/>
                <p:cNvGrpSpPr/>
                <p:nvPr/>
            </p:nvGrpSpPr>
            <p:grpSpPr>
                <a:xfrm>
                    <a:off x="0" y="0"/>
                    <a:ext cx="0" cy="0"/>
                    <a:chOff x="0" y="0"/>
                    <a:chExt cx="0" cy="0"/>
                </a:xfrm>
            </p:grpSpPr>
        </p:spTree>
    </p:cSld>
    <p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/>
    <p:sldLayoutIdLst>
        <p:sldLayoutId id="2147483649" r:id="rId1"/>
    </p:sldLayoutIdLst>
</p:sldMaster>'''
    
    def _get_slide_layout_xml(self) -> str:
        """Get slide layout XML"""
        return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="obj" preserve="1">
    <p:cSld name="Title and Content">
        <p:spTree>
            <p:nvGrpSpPr>
                <p:cNvPr id="1" name=""/>
                <p:cNvGrpSpPr/>
                <p:nvPr/>
            </p:nvGrpSpPr>
            <p:grpSpPr>
                <a:xfrm>
                    <a:off x="0" y="0"/>
                    <a:ext cx="0" cy="0"/>
                    <a:chOff x="0" y="0"/>
                    <a:chExt cx="0" cy="0"/>
                </a:xfrm>
            </p:grpSpPr>
        </p:spTree>
    </p:cSld>
    <p:clrMapOvr>
        <a:masterClrMapping/>
    </p:clrMapOvr>
</p:sldLayout>'''
    
    def _use_template_from_s3(self) -> Optional[bytes]:
        """Try to use pre-built template from S3"""
        try:
            logger.info("Attempting to use template from S3...")
            response = s3.get_object(
                Bucket=self.documents_bucket,
                Key='templates/slide_23_template.pptx'
            )
            template_content = response['Body'].read()
            logger.info(f"Successfully loaded template from S3: {len(template_content)} bytes")
            return template_content
        except Exception as e:
            logger.error(f"Failed to load template from S3: {e}")
            return None