"""
Full Slide Extractor
Extracts a single slide while preserving all content
"""

import io
import logging
from typing import Optional
from pptx import Presentation

logger = logging.getLogger(__name__)

class FullSlideExtractor:
    def extract_single_slide(self, pptx_bytes: bytes, slide_position: int) -> bytes:
        """
        Extract a single slide from presentation by removing all others
        
        Args:
            pptx_bytes: The PowerPoint file as bytes
            slide_position: The 1-based position of the slide to keep
        
        Returns:
            PowerPoint with only the requested slide
        """
        
        # Load presentation
        prs = Presentation(io.BytesIO(pptx_bytes))
        
        # Validate slide position
        if slide_position < 1 or slide_position > len(prs.slides):
            raise ValueError(f"Invalid slide position: {slide_position}")
        
        logger.info(f"Extracting slide at position {slide_position} from {len(prs.slides)} slides")
        
        # Remove all slides except the target one
        # Work backwards to avoid index issues
        slides_to_remove = []
        for i in range(len(prs.slides)):
            if i != slide_position - 1:  # Keep the target slide
                slides_to_remove.append(i)
        
        # Remove slides in reverse order
        for idx in reversed(slides_to_remove):
            # Get the slide element
            slide_element = prs.slides._sldIdLst[idx]
            slide_id = slide_element.rId
            
            # Remove from slide ID list
            prs.slides._sldIdLst.remove(slide_element)
            
            # Remove the slide part
            if hasattr(prs.part, 'drop_rel'):
                prs.part.drop_rel(slide_id)
            elif hasattr(prs.part, '_rels'):
                # Alternative method
                if slide_id in prs.part._rels:
                    del prs.part._rels[slide_id]
            
            logger.info(f"Removed slide at position {idx + 1}")
        
        logger.info(f"Kept slide at position {slide_position}")
        
        # Save the single-slide presentation
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return output.getvalue()


def extract_single_slide_full(template_bytes: bytes, slide_number: int) -> bytes:
    """Helper function for full extraction"""
    extractor = FullSlideExtractor()
    return extractor.extract_single_slide(template_bytes, slide_number)