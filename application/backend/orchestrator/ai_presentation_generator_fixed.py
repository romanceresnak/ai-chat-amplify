"""
AI-Powered Presentation Generator for Financial Services - Fixed Version
Generates professional presentations using Bedrock AI and python-pptx
"""

import json
import boto3
from typing import Dict, List, Any
import logging
import io
import re

logger = logging.getLogger()
logger.setLevel(logging.INFO)

class AIPresentationGenerator:
    def __init__(self):
        self.bedrock_runtime = boto3.client('bedrock-runtime', region_name='us-east-1')
        self.model_id = 'eu.anthropic.claude-3-5-sonnet-20240620-v1:0'
    
    def analyze_presentation_request(self, instructions: str) -> Dict[str, Any]:
        """Use AI to analyze and structure the presentation request"""
        
        prompt = f"""Analyze this presentation request and provide a structured plan.
        
Request: {instructions}

Return a JSON structure with:
1. presentation_type: "financial_pe", "financial_ib", "loan_portfolio", "general"
2. title: Main presentation title
3. num_slides: Number of slides to create
4. sections: Array of sections, each containing:
   - title: Section title
   - slide_type: "title", "content", "chart"
   - content: Key points or data
   - chart_data: If applicable, data for charts

Be specific and detailed."""

        request_body = {
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 2000,
            "temperature": 0.7,
            "messages": [
                {
                    "role": "user",
                    "content": f"{prompt}\n\nReturn only valid JSON in your response."
                }
            ]
        }
        
        try:
            response = self.bedrock_runtime.invoke_model(
                modelId=self.model_id,
                body=json.dumps(request_body),
                contentType='application/json'
            )
            
            response_body = json.loads(response['body'].read())
            response_text = response_body.get('content', [{}])[0].get('text', '')
            
            # Extract JSON from response
            json_match = re.search(r'```json\n(.*?)\n```', response_text, re.DOTALL)
            if json_match:
                return json.loads(json_match.group(1))
        except Exception as e:
            logger.error(f"Error analyzing request: {e}")
        
        # Fallback structure
        return self._get_default_structure(instructions)
    
    def _get_default_structure(self, instructions: str) -> Dict[str, Any]:
        """Provide default structure based on keywords"""
        
        instructions_lower = instructions.lower()
        
        if "loan portfolio" in instructions_lower:
            return self._get_loan_portfolio_structure()
        else:
            return self._get_general_financial_structure()
    
    def _get_loan_portfolio_structure(self) -> Dict[str, Any]:
        """Loan portfolio presentation structure"""
        return {
            "presentation_type": "loan_portfolio",
            "title": "Loan Portfolio Analysis",
            "num_slides": 10,
            "sections": [
                {
                    "title": "South Plains Financial, Inc.",
                    "slide_type": "title",
                    "content": ["September 2024", "Loan Portfolio Analysis", "Investor Presentation"]
                },
                {
                    "title": "Loan Portfolio Overview",
                    "slide_type": "chart",
                    "content": ["Portfolio composition", "Net Loans: $2.3 Billion"],
                    "chart_data": {
                        "type": "doughnut",
                        "categories": [
                            "Commercial Real Estate",
                            "Commercial – General", 
                            "Commercial – Specialized",
                            "1–4 Family Residential",
                            "Auto Loans"
                        ],
                        "values": [28, 27, 14, 15, 16]
                    }
                },
                {
                    "title": "Portfolio Highlights",
                    "slide_type": "content",
                    "content": [
                        "Commercial Real Estate: Well-diversified across property types",
                        "Commercial – General: Strong relationships with local businesses",
                        "Commercial – Specialized: Strategic focus on agricultural and energy",
                        "Residential: Conservative underwriting standards",
                        "Auto Loans: Indirect lending through dealer network"
                    ]
                },
                {
                    "title": "Credit Quality Metrics",
                    "slide_type": "content",
                    "content": [
                        "Non-Performing Loans: 0.6% of total loans",
                        "Allowance for Credit Losses: 1.2% of loans",
                        "Net Charge-offs: 0.15% YTD",
                        "Strong underwriting and risk management"
                    ]
                },
                {
                    "title": "Geographic Distribution",
                    "slide_type": "content",
                    "content": [
                        "Texas Panhandle: 65%",
                        "West Texas: 20%",
                        "Eastern New Mexico: 10%",
                        "Other Markets: 5%"
                    ]
                }
            ]
        }
    
    def _get_general_financial_structure(self) -> Dict[str, Any]:
        """General financial presentation structure"""
        return {
            "presentation_type": "general",
            "title": "Financial Presentation",
            "num_slides": 5,
            "sections": [
                {
                    "title": "Financial Analysis",
                    "slide_type": "title",
                    "content": ["Executive Presentation"]
                },
                {
                    "title": "Overview",
                    "slide_type": "content",
                    "content": ["Key highlights", "Financial metrics", "Strategic initiatives"]
                }
            ]
        }
    
    def generate_presentation(self, instructions: str) -> bytes:
        """Generate a complete presentation based on instructions"""
        
        try:
            # Import python-pptx dynamically
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.chart.data import ChartData
            from pptx.enum.chart import XL_CHART_TYPE
            
            logger.info("Successfully imported python-pptx")
            
        except ImportError as e:
            logger.error(f"Failed to import python-pptx: {e}")
            raise ImportError("python-pptx is not available in the Lambda environment")
        
        # Analyze the request
        structure = self.analyze_presentation_request(instructions)
        logger.info(f"Generated structure: {json.dumps(structure)}")
        
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Generate slides based on structure
        for section in structure['sections']:
            try:
                if section['slide_type'] == 'title':
                    self._create_title_slide(prs, section)
                elif section['slide_type'] == 'content':
                    self._create_content_slide(prs, section)
                elif section['slide_type'] == 'chart':
                    self._create_chart_slide(prs, section, ChartData, XL_CHART_TYPE, Inches, Pt, RGBColor)
            except Exception as e:
                logger.error(f"Error creating slide: {e}")
                # Add fallback content slide
                self._create_content_slide(prs, section)
        
        # Save to bytes using getvalue() like presentation_agent.py
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_content = pptx_bytes.getvalue()  # This is the key fix!
        
        logger.info(f"Generated presentation with {len(prs.slides)} slides, size: {len(pptx_content)} bytes")
        
        return pptx_content
    
    def _create_title_slide(self, prs, section: Dict):
        """Create a title slide using standard layout"""
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = section.get('title', 'Presentation')
        
        # Set subtitle if available
        if section.get('content') and len(slide.placeholders) > 1:
            slide.placeholders[1].text = '\n'.join(section['content'])
    
    def _create_content_slide(self, prs, section: Dict):
        """Create a content slide using standard layout"""
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = section.get('title', 'Content')
        
        # Set content
        try:
            content_placeholder = slide.placeholders[1]
            if hasattr(content_placeholder, 'text_frame'):
                tf = content_placeholder.text_frame
                
                for i, point in enumerate(section.get('content', [])):
                    if i == 0:
                        tf.text = point
                    else:
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0
        except Exception as e:
            logger.warning(f"Could not add content to slide: {e}")
    
    def _create_chart_slide(self, prs, section: Dict, ChartData, XL_CHART_TYPE, Inches, Pt, RGBColor):
        """Create a chart slide"""
        slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[-1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = section.get('title', 'Chart')
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        
        # Add chart if data provided
        chart_data = section.get('chart_data', {})
        if chart_data.get('type') == 'doughnut' and chart_data.get('categories'):
            try:
                # Create chart data
                chart_data_obj = ChartData()
                chart_data_obj.categories = chart_data['categories']
                chart_data_obj.add_series('Portfolio', chart_data['values'])
                
                # Add chart
                x, y, cx, cy = Inches(1), Inches(2), Inches(6), Inches(5)
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data_obj
                ).chart
                
                # Style the chart
                chart.has_legend = True
                
            except Exception as e:
                logger.error(f"Error creating chart: {e}")
                # Add content as fallback
                self._create_content_slide(prs, section)