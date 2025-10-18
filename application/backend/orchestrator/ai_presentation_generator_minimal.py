"""
Minimal AI-Powered Presentation Generator - Debug Version
Creates very basic presentations to isolate corruption issues
"""

import json
import boto3
from typing import Dict, Any
import logging
import io

logger = logging.getLogger()

# Global variables for python-pptx availability
PPTX_AVAILABLE = False
Presentation = None
Inches = None
Pt = None

class AIPresentationGenerator:
    def __init__(self):
        self.bedrock_runtime = boto3.client('bedrock-runtime', region_name='us-east-1')
        self.model_id = 'eu.anthropic.claude-3-5-sonnet-20240620-v1:0'
        self._initialize_pptx()
    
    def _initialize_pptx(self):
        """Initialize python-pptx imports on demand"""
        global PPTX_AVAILABLE, Presentation, Inches, Pt
        
        if PPTX_AVAILABLE:
            return
        
        try:
            import sys
            import os
            
            layer_path = "/opt/python/lib/python3.11/site-packages"
            if os.path.exists(layer_path) and layer_path not in sys.path:
                sys.path.insert(0, layer_path)
            
            from pptx import Presentation as _Presentation
            from pptx.util import Inches as _Inches, Pt as _Pt
            
            Presentation = _Presentation
            Inches = _Inches
            Pt = _Pt
            PPTX_AVAILABLE = True
            
            logger.info("python-pptx successfully imported")
            
        except ImportError as e:
            logger.error(f"Failed to import python-pptx: {e}")
            PPTX_AVAILABLE = False
    
    def analyze_presentation_request(self, instructions: str) -> Dict[str, Any]:
        """Return minimal structure for testing"""
        return self._get_loan_portfolio_structure()
    
    def _get_loan_portfolio_structure(self) -> Dict[str, Any]:
        """Minimal loan portfolio structure"""
        return {
            "presentation_type": "loan_portfolio",
            "title": "Loan Portfolio Analysis",
            "sections": [
                {
                    "title": "South Plains Financial, Inc.",
                    "slide_type": "title",
                    "content": ["Loan Portfolio Analysis", "September 2024"]
                },
                {
                    "title": "Portfolio Overview",
                    "slide_type": "content",
                    "content": [
                        "Total Loans: $2.3 Billion",
                        "Commercial Real Estate: 28%",
                        "Commercial General: 27%",
                        "Residential: 15%"
                    ]
                },
                {
                    "title": "Key Highlights",
                    "slide_type": "content",
                    "content": [
                        "Strong asset quality",
                        "Diversified portfolio",
                        "Conservative underwriting"
                    ]
                }
            ]
        }
    
    def generate_presentation(self, instructions: str) -> bytes:
        """Generate minimal presentation"""
        
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not available")
        
        try:
            # Create new presentation
            prs = Presentation()
            
            # Get structure
            structure = self.analyze_presentation_request(instructions)
            
            # Create slides
            for section in structure['sections']:
                if section['slide_type'] == 'title':
                    self._create_title_slide(prs, section)
                else:
                    self._create_content_slide(prs, section)
            
            # Save to bytes
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            
            return output.read()
            
        except Exception as e:
            logger.error(f"Error generating presentation: {e}")
            raise
    
    def _create_title_slide(self, prs, section: Dict):
        """Create minimal title slide"""
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = section.get('title', 'Title')
        
        # Set subtitle if exists
        if len(slide.placeholders) > 1 and section.get('content'):
            subtitle = slide.placeholders[1]
            subtitle.text = '\n'.join(section['content'])
    
    def _create_content_slide(self, prs, section: Dict):
        """Create minimal content slide"""
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = section.get('title', 'Content')
        
        # Set content
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            if hasattr(content, 'text_frame'):
                tf = content.text_frame
                
                # Add bullet points
                for i, point in enumerate(section.get('content', [])):
                    if i == 0:
                        tf.text = point
                    else:
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0