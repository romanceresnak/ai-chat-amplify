"""
AI-Powered Presentation Generator for Financial Services
Generates professional presentations using Bedrock AI and python-pptx
"""

import json
import boto3
from typing import Dict, List, Any, Optional
import logging
import io
import re

logger = logging.getLogger()

# Global variables for python-pptx availability
PPTX_AVAILABLE = False
Presentation = None
Inches = None
Pt = None
RGBColor = None
PP_ALIGN = None
ChartData = None
XL_CHART_TYPE = None
MSO_SHAPE = None

class AIPresentationGenerator:
    def __init__(self):
        self.bedrock_runtime = boto3.client('bedrock-runtime', region_name='eu-west-1')
        self.model_id = 'eu.anthropic.claude-3-5-sonnet-20240620-v1:0'
        self._initialize_pptx()
    
    def _initialize_pptx(self):
        """Initialize python-pptx imports on demand"""
        global PPTX_AVAILABLE, Presentation, Inches, Pt, RGBColor, PP_ALIGN, ChartData, XL_CHART_TYPE, MSO_SHAPE
        
        if PPTX_AVAILABLE:
            return  # Already initialized
        
        try:
            import sys
            import os
            
            # Debug info
            logger.info(f"Python path: {sys.path}")
            layer_path = "/opt/python/lib/python3.11/site-packages"
            
            if os.path.exists(layer_path):
                logger.info(f"Layer directory exists: {layer_path}")
                if layer_path not in sys.path:
                    sys.path.insert(0, layer_path)
                    logger.info(f"Added {layer_path} to sys.path")
            
            # Import python-pptx components
            from pptx import Presentation as _Presentation
            from pptx.util import Inches as _Inches, Pt as _Pt
            from pptx.dml.color import RGBColor as _RGBColor
            from pptx.enum.text import PP_ALIGN as _PP_ALIGN
            from pptx.chart.data import ChartData as _ChartData
            from pptx.enum.chart import XL_CHART_TYPE as _XL_CHART_TYPE
            from pptx.enum.shapes import MSO_SHAPE as _MSO_SHAPE
            
            # Set global variables
            Presentation = _Presentation
            Inches = _Inches
            Pt = _Pt
            RGBColor = _RGBColor
            PP_ALIGN = _PP_ALIGN
            ChartData = _ChartData
            XL_CHART_TYPE = _XL_CHART_TYPE
            MSO_SHAPE = _MSO_SHAPE
            PPTX_AVAILABLE = True
            
            logger.info("python-pptx successfully imported in AI generator")
            
        except ImportError as e:
            logger.error(f"Failed to import python-pptx: {e}")
            PPTX_AVAILABLE = False
        
    def analyze_presentation_request(self, instructions: str) -> Dict[str, Any]:
        """Use AI to analyze and structure the presentation request"""
        
        prompt = f"""Analyze this presentation request and provide a structured plan.
        
Request: {instructions}

Return a JSON structure with:
1. presentation_type: "financial_pe", "financial_ib", "loan_portfolio", "general"
2. title: Main presentation title
3. num_slides: Number of slides to create (default 40 for PE/IB decks)
4. sections: Array of sections, each containing:
   - title: Section title
   - slide_type: "title", "content", "chart", "table", "mixed"
   - content: Key points or data
   - chart_data: If applicable, data for charts
5. design_theme: "corporate", "financial", "modern"
6. color_scheme: Primary colors to use

For financial presentations, include standard sections like:
- Executive Summary
- Company Overview
- Financial Performance
- Market Analysis
- Investment Thesis
- Risk Factors
- Appendices

Be specific and detailed."""

        # Use the correct format for Claude 3.5 Sonnet
        request_body = {
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 2000,
            "temperature": 0.7,
            "messages": [
                {
                    "role": "user",
                    "content": f"{prompt}\n\nReturn only valid JSON in your response."
                }
            ]
        }
        
        response = self.bedrock_runtime.invoke_model(
            modelId=self.model_id,
            body=json.dumps(request_body),
            contentType='application/json'
        )
        
        response_body = json.loads(response['body'].read())
        response_text = response_body.get('content', [{}])[0].get('text', '')
        
        # Extract JSON from response
        json_match = re.search(r'```json\n(.*?)\n```', response_text, re.DOTALL)
        if json_match:
            return json.loads(json_match.group(1))
        else:
            # Fallback structure
            return self._get_default_structure(instructions)
    
    def _get_default_structure(self, instructions: str) -> Dict[str, Any]:
        """Provide default structure based on keywords"""
        
        instructions_lower = instructions.lower()
        
        if "private equity" in instructions_lower or "investment committee" in instructions_lower:
            return self._get_pe_deck_structure()
        elif "debt issuance" in instructions_lower:
            return self._get_debt_issuance_structure()
        elif "loan portfolio" in instructions_lower:
            return self._get_loan_portfolio_structure()
        else:
            return self._get_general_financial_structure()
    
    def _get_pe_deck_structure(self) -> Dict[str, Any]:
        """Standard PE investment committee deck structure"""
        return {
            "presentation_type": "financial_pe",
            "title": "Private Equity Investment Committee Presentation",
            "num_slides": 40,
            "design_theme": "financial",
            "color_scheme": {"primary": "002855", "accent": "CC0000"},
            "sections": [
                {
                    "title": "Title Slide",
                    "slide_type": "title",
                    "content": ["Investment Committee Meeting", "Confidential"]
                },
                {
                    "title": "Executive Summary",
                    "slide_type": "content",
                    "content": [
                        "Investment opportunity overview",
                        "Key investment highlights",
                        "Expected returns and timeline",
                        "Recommendation"
                    ]
                },
                {
                    "title": "Investment Thesis",
                    "slide_type": "mixed",
                    "content": [
                        "Strategic rationale",
                        "Value creation opportunities",
                        "Competitive advantages",
                        "Market positioning"
                    ]
                },
                {
                    "title": "Company Overview",
                    "slide_type": "content",
                    "content": [
                        "Business description",
                        "Products and services",
                        "Management team",
                        "Corporate structure"
                    ]
                },
                {
                    "title": "Financial Performance",
                    "slide_type": "chart",
                    "content": ["Historical financials", "Revenue growth", "EBITDA margins"],
                    "chart_data": {
                        "type": "column",
                        "categories": ["2020", "2021", "2022", "2023"],
                        "series": [
                            {"name": "Revenue", "values": [100, 120, 145, 170]},
                            {"name": "EBITDA", "values": [20, 28, 35, 42]}
                        ]
                    }
                },
                {
                    "title": "Market Analysis",
                    "slide_type": "mixed",
                    "content": [
                        "Market size and growth",
                        "Competitive landscape",
                        "Market trends",
                        "Growth drivers"
                    ]
                },
                {
                    "title": "Transaction Structure",
                    "slide_type": "table",
                    "content": [
                        "Deal terms",
                        "Valuation metrics",
                        "Sources and uses",
                        "Capital structure"
                    ]
                },
                {
                    "title": "Value Creation Plan",
                    "slide_type": "content",
                    "content": [
                        "Operational improvements",
                        "Revenue growth initiatives",
                        "Cost optimization",
                        "Strategic acquisitions"
                    ]
                },
                {
                    "title": "Risk Analysis",
                    "slide_type": "content",
                    "content": [
                        "Key risks and mitigants",
                        "Market risks",
                        "Operational risks",
                        "Financial risks"
                    ]
                },
                {
                    "title": "Exit Strategy",
                    "slide_type": "content",
                    "content": [
                        "Exit timing",
                        "Exit options",
                        "Expected returns",
                        "Value realization"
                    ]
                }
            ]
        }
    
    def _get_debt_issuance_structure(self) -> Dict[str, Any]:
        """Debt issuance presentation structure"""
        return {
            "presentation_type": "financial_ib",
            "title": "Debt Issuance Presentation",
            "num_slides": 35,
            "design_theme": "corporate",
            "color_scheme": {"primary": "003366", "accent": "FF6600"},
            "sections": [
                {
                    "title": "Title Slide",
                    "slide_type": "title",
                    "content": ["Debt Issuance Opportunity", "Institutional Investor Presentation"]
                },
                {
                    "title": "Transaction Overview",
                    "slide_type": "content",
                    "content": [
                        "Issuer overview",
                        "Transaction size and structure",
                        "Use of proceeds",
                        "Key terms"
                    ]
                },
                {
                    "title": "Company Profile",
                    "slide_type": "mixed",
                    "content": [
                        "Business overview",
                        "Market position",
                        "Competitive advantages",
                        "Management team"
                    ]
                },
                {
                    "title": "Financial Highlights",
                    "slide_type": "chart",
                    "content": ["Revenue trends", "Profitability metrics", "Cash flow generation"],
                    "chart_data": {
                        "type": "line",
                        "categories": ["Q1'23", "Q2'23", "Q3'23", "Q4'23"],
                        "series": [
                            {"name": "Revenue", "values": [250, 265, 280, 295]},
                            {"name": "EBITDA", "values": [50, 55, 58, 62]}
                        ]
                    }
                },
                {
                    "title": "Credit Profile",
                    "slide_type": "table",
                    "content": [
                        "Credit ratings",
                        "Leverage metrics",
                        "Coverage ratios",
                        "Debt maturity profile"
                    ]
                },
                {
                    "title": "Investment Highlights",
                    "slide_type": "content",
                    "content": [
                        "Strong market position",
                        "Stable cash flows",
                        "Conservative leverage",
                        "Experienced management"
                    ]
                }
            ]
        }
    
    def _get_loan_portfolio_structure(self) -> Dict[str, Any]:
        """Loan portfolio presentation structure with professional design"""
        return {
            "presentation_type": "loan_portfolio",
            "title": "Loan Portfolio Analysis",
            "num_slides": 15,
            "design_theme": "modern",
            "color_scheme": {"primary": "1F497D", "accent": "C0504D"},
            "sections": [
                {
                    "title": "South Plains Financial, Inc.",
                    "slide_type": "title",
                    "content": ["September 2024", "Loan Portfolio Analysis", "Investor Presentation"]
                },
                {
                    "title": "Loan Portfolio",
                    "slide_type": "chart",
                    "content": ["Portfolio composition overview", "Diversified loan book across multiple segments"],
                    "chart_data": {
                        "type": "doughnut",
                        "categories": [
                            "Commercial Real Estate",
                            "Commercial – General", 
                            "Commercial – Specialized",
                            "1–4 Family Residential",
                            "Auto Loans",
                            "Construction",
                            "Other Consumer"
                        ],
                        "values": [28, 27, 14, 15, 9, 4, 3],
                        "center_text": "Net Loans\n2Q'24: $2.3 Billion"
                    }
                },
                {
                    "title": "Portfolio Highlights",
                    "slide_type": "content",
                    "content": [
                        "Commercial Real Estate: Well-diversified across property types with focus on income-producing assets",
                        "Commercial – General: Strong relationships with local businesses, including SBA and PPP lending",
                        "Commercial – Specialized: Strategic focus on agricultural and energy sectors given regional expertise",
                        "1–4 Family Residential: Conservative underwriting with focus on owner-occupied properties",
                        "Auto Loans: Indirect lending through established dealer network"
                    ]
                },
                {
                    "title": "Financial Projections",
                    "slide_type": "chart", 
                    "content": ["Loan growth outlook", "Credit quality trends"],
                    "chart_data": {
                        "type": "column",
                        "categories": ["2021", "2022", "2023", "2024E"],
                        "series": [
                            {"name": "Total Loans ($B)", "values": [2.0, 2.1, 2.2, 2.3]},
                            {"name": "NPL Ratio (%)", "values": [0.8, 0.9, 0.7, 0.6]}
                        ]
                    }
                },
                {
                    "title": "Due Diligence Findings",
                    "slide_type": "content",
                    "content": [
                        "Strong underwriting standards maintained throughout credit cycle",
                        "Experienced credit team with deep local market knowledge", 
                        "Conservative loan-to-value ratios across all asset classes",
                        "Robust portfolio monitoring and early warning systems"
                    ]
                },
                {
                    "title": "Comparative Analysis",
                    "slide_type": "content",
                    "content": [
                        "Peer-leading asset quality metrics",
                        "Competitive net interest margins",
                        "Above-average loan growth in target markets",
                        "Strong deposit franchise supporting funding costs"
                    ]
                },
                {
                    "title": "Pathway Analysis",
                    "slide_type": "content", 
                    "content": [
                        "Continued focus on relationship-based commercial lending",
                        "Selective expansion in adjacent markets",
                        "Technology investments to improve efficiency",
                        "Maintain disciplined credit culture"
                    ]
                },
                {
                    "title": "Integration Plan",
                    "slide_type": "content",
                    "content": [
                        "100-day integration roadmap established",
                        "Key personnel retention agreements in place",
                        "Systems integration timeline defined",
                        "Customer communication strategy developed"
                    ]
                },
                {
                    "title": "Financing Structure",
                    "slide_type": "content",
                    "content": [
                        "Transaction structure and terms",
                        "Sources and uses of funds",
                        "Pro forma capital ratios",
                        "Expected synergies and cost savings"
                    ]
                },
                {
                    "title": "Regulatory Considerations",
                    "slide_type": "content",
                    "content": [
                        "Regulatory approval timeline",
                        "Required filings and documentation", 
                        "Community Reinvestment Act commitments",
                        "Capital and liquidity requirements"
                    ]
                }
            ]
        }
    
    def _get_general_financial_structure(self) -> Dict[str, Any]:
        """General financial presentation structure"""
        return {
            "presentation_type": "general",
            "title": "Financial Presentation",
            "num_slides": 10,
            "design_theme": "corporate",
            "color_scheme": {"primary": "1F2937", "accent": "3B82F6"},
            "sections": [
                {
                    "title": "Title Slide",
                    "slide_type": "title",
                    "content": ["Financial Analysis", "Executive Presentation"]
                },
                {
                    "title": "Overview",
                    "slide_type": "content",
                    "content": ["Key highlights", "Financial metrics", "Strategic initiatives"]
                }
            ]
        }
    
    def generate_presentation(self, instructions: str) -> bytes:
        """Generate a complete presentation based on instructions"""
        
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not available in the Lambda environment")
        
        # Analyze the request
        structure = self.analyze_presentation_request(instructions)
        
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Generate slides based on structure
        for i, section in enumerate(structure['sections']):
            if section['slide_type'] == 'title':
                self._create_title_slide(prs, section, structure['color_scheme'])
            elif section['slide_type'] == 'content':
                self._create_content_slide(prs, section, structure['color_scheme'])
            elif section['slide_type'] == 'chart':
                self._create_chart_slide(prs, section, structure['color_scheme'])
            elif section['slide_type'] == 'table':
                self._create_table_slide(prs, section, structure['color_scheme'])
            elif section['slide_type'] == 'mixed':
                self._create_mixed_slide(prs, section, structure['color_scheme'])
        
        # Fill remaining slides if needed for PE/IB decks
        current_slides = len(prs.slides)
        target_slides = structure.get('num_slides', 10)
        
        if current_slides < target_slides:
            self._fill_remaining_slides(prs, structure, current_slides, target_slides)
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return output.read()
    
    def _create_title_slide(self, prs: Presentation, section: Dict, colors: Dict):
        """Create a professional title slide with modern design"""
        # Use blank layout (usually index 5 or 6, fallback to last available)
        try:
            slide_layout = prs.slide_layouts[5]  # Try standard blank layout
        except IndexError:
            slide_layout = prs.slide_layouts[-1]  # Use last available layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(51, 63, 72)  # Dark sophisticated gray
        
        # Add company logo placeholder (top-left corner)
        logo_shape = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(2), Inches(0.8))
        logo_frame = logo_shape.text_frame
        logo_frame.text = "South Plains Financial, Inc."
        logo_frame.paragraphs[0].font.size = Pt(18)
        logo_frame.paragraphs[0].font.bold = True
        logo_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Add main title (center)
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
        title_frame = title_shape.text_frame
        title_frame.text = section.get('title', 'Financial Presentation')
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].font.name = 'Arial'
        
        # Add subtitle
        if section.get('content'):
            subtitle_shape = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.text = '\n'.join(section['content'])
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
            subtitle_frame.paragraphs[0].font.name = 'Arial'
        
        # Add date and confidentiality notice (bottom)
        date_shape = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.333), Inches(0.8))
        date_frame = date_shape.text_frame
        date_frame.text = "September 2024 • CONFIDENTIAL"
        date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        date_frame.paragraphs[0].font.size = Pt(16)
        date_frame.paragraphs[0].font.color.rgb = RGBColor(160, 160, 160)
        date_frame.paragraphs[0].font.name = 'Arial'
        
        # Add decorative accent line as a thin rectangle
        accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.8), Inches(5.333), Pt(3))
        accent_fill = accent_line.fill
        accent_fill.solid()
        accent_fill.fore_color.rgb = RGBColor(192, 80, 77)  # South Plains red accent
        accent_line.line.fill.background()  # No outline
    
    def _create_content_slide(self, prs: Presentation, section: Dict, colors: Dict):
        """Create a professional content slide with modern design"""
        # Use blank layout safely
        try:
            slide_layout = prs.slide_layouts[5]  # Try standard blank layout
        except IndexError:
            slide_layout = prs.slide_layouts[-1]  # Use last available layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set professional light background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 250, 252)  # Very light gray
        
        # Add header bar as a rectangle shape instead of textbox
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
        header_fill = header_shape.fill
        header_fill.solid()
        header_fill.fore_color.rgb = RGBColor(31, 73, 125)  # South Plains blue
        header_shape.line.fill.background()  # No outline
        
        # Add title in header
        title_shape = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.333), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = section.get('title', 'Content')
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].font.name = 'Arial'
        
        # Add content with better formatting
        content_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(5))
        content_frame = content_shape.text_frame
        content_frame.margin_left = Inches(0.3)
        content_frame.margin_top = Inches(0.2)
        
        for i, point in enumerate(section.get('content', [])):
            if i == 0:
                content_frame.text = f"• {point}"
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
                p.text = f"• {point}"
            
            # Style each bullet point
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(55, 65, 81)  # Dark gray
            p.font.name = 'Arial'
            p.space_after = Pt(12)  # Add space between points
            
        # Add accent line at bottom as a rectangle
        accent_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(7), Inches(2), Inches(0.1))
        accent_fill = accent_shape.fill
        accent_fill.solid()
        accent_fill.fore_color.rgb = RGBColor(192, 80, 77)  # South Plains red accent
        accent_shape.line.fill.background()  # No outline
    
    def _create_chart_slide(self, prs: Presentation, section: Dict, colors: Dict):
        """Create a professional slide with charts"""
        # Use blank layout safely
        try:
            slide_layout = prs.slide_layouts[5]  # Try standard blank layout
        except IndexError:
            slide_layout = prs.slide_layouts[-1]  # Use last available layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set professional light background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 250, 252)  # Very light gray
        
        # Add header bar as a rectangle
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
        header_fill = header_shape.fill
        header_fill.solid()
        header_fill.fore_color.rgb = RGBColor(31, 73, 125)  # South Plains blue
        header_shape.line.fill.background()  # No outline
        
        # Add title in header
        title_shape = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.333), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = section.get('title', 'Analysis')
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].font.name = 'Arial'
        
        # Add chart based on data
        chart_data = section.get('chart_data', {})
        
        if chart_data.get('type') == 'doughnut':
            self._create_doughnut_chart(slide, chart_data, colors)
        elif chart_data.get('type') == 'column':
            self._create_column_chart(slide, chart_data, colors)
        elif chart_data.get('type') == 'line':
            self._create_line_chart(slide, chart_data, colors)
    
    def _create_doughnut_chart(self, slide, chart_data: Dict, colors: Dict):
        """Create a professional doughnut chart with modern styling"""
        chart_data_obj = ChartData()
        chart_data_obj.categories = chart_data['categories']
        chart_data_obj.add_series('Portfolio', chart_data['values'])
        
        # Position chart on left side with better spacing
        x, y, cx, cy = Inches(0.8), Inches(2), Inches(6), Inches(4.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data_obj
        ).chart
        
        # Style the chart
        chart.has_legend = True
        chart.legend.position = 9  # Right position
        chart.legend.font.size = Pt(12)
        
        # Add center text if provided with better styling
        if 'center_text' in chart_data:
            center_x = x + cx/2 - Inches(1.8)
            center_y = y + cy/2 - Inches(0.4)
            textbox = slide.shapes.add_textbox(center_x, center_y, Inches(3.6), Inches(0.8))
            text_frame = textbox.text_frame
            text_frame.text = chart_data['center_text']
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
            text_frame.paragraphs[0].font.name = 'Arial'
        
        # Add professional legend/highlights on the right
        if chart_data.get('categories'):
            legend_x = Inches(7.5)
            legend_y = Inches(2.5)
            
            # Create legend header
            header_shape = slide.shapes.add_textbox(legend_x, legend_y - Inches(0.5), Inches(5), Inches(0.4))
            header_frame = header_shape.text_frame
            header_frame.text = "Portfolio Composition"
            header_frame.paragraphs[0].font.size = Pt(16)
            header_frame.paragraphs[0].font.bold = True
            header_frame.paragraphs[0].font.color.rgb = RGBColor(55, 65, 81)
            header_frame.paragraphs[0].font.name = 'Arial'
            
            # Color palette for legend items
            legend_colors = [
                RGBColor(79, 129, 189),   # Theme blue
                RGBColor(192, 80, 77),    # Theme red  
                RGBColor(155, 187, 89),   # Theme green
                RGBColor(128, 100, 162),  # Theme purple
                RGBColor(247, 150, 70),   # Theme orange
                RGBColor(75, 172, 198),   # Theme teal
                RGBColor(156, 163, 175),  # Gray
            ]
            
            for i, (category, value) in enumerate(zip(chart_data['categories'], chart_data['values'])):
                y_pos = legend_y + (i * Inches(0.5))
                
                # Color indicator as a small rectangle
                color_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, legend_x, y_pos, Inches(0.3), Inches(0.3))
                color_fill = color_shape.fill
                color_fill.solid()
                color_fill.fore_color.rgb = legend_colors[i % len(legend_colors)]
                color_shape.line.fill.background()  # No outline
                
                # Category label with percentage
                label_shape = slide.shapes.add_textbox(legend_x + Inches(0.4), y_pos, Inches(4.5), Inches(0.4))
                label_frame = label_shape.text_frame
                label_frame.text = f"{category}: {value}%"
                label_frame.paragraphs[0].font.size = Pt(14)
                label_frame.paragraphs[0].font.color.rgb = RGBColor(55, 65, 81)
                label_frame.paragraphs[0].font.name = 'Arial'
    
    def _create_column_chart(self, slide, chart_data: Dict, colors: Dict):
        """Create a column chart"""
        chart_data_obj = ChartData()
        chart_data_obj.categories = chart_data['categories']
        
        for series in chart_data['series']:
            chart_data_obj.add_series(series['name'], series['values'])
        
        x, y, cx, cy = Inches(1.5), Inches(2), Inches(10), Inches(5)
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_obj
        )
    
    def _create_line_chart(self, slide, chart_data: Dict, colors: Dict):
        """Create a line chart"""
        chart_data_obj = ChartData()
        chart_data_obj.categories = chart_data['categories']
        
        for series in chart_data['series']:
            chart_data_obj.add_series(series['name'], series['values'])
        
        x, y, cx, cy = Inches(1.5), Inches(2), Inches(10), Inches(5)
        slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data_obj
        )
    
    def _create_table_slide(self, prs: Presentation, section: Dict, colors: Dict):
        """Create a slide with a table"""
        # Use blank layout safely
        try:
            slide_layout = prs.slide_layouts[5]  # Try standard blank layout
        except IndexError:
            slide_layout = prs.slide_layouts[-1]  # Use last available layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = section.get('title', 'Table')
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        
        # Create table (simplified for now)
        rows = len(section.get('content', [])) + 1
        cols = 2
        left = Inches(1)
        top = Inches(2)
        width = Inches(11)
        height = Inches(0.8 * rows)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Add headers
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        
        # Add content
        for i, item in enumerate(section.get('content', [])):
            table.cell(i+1, 0).text = item
            table.cell(i+1, 1).text = "TBD"  # Would be filled with actual data
        
        self._apply_colors(slide, colors)
    
    def _create_mixed_slide(self, prs: Presentation, section: Dict, colors: Dict):
        """Create a slide with mixed content (text and visuals)"""
        # Use blank layout safely
        try:
            slide_layout = prs.slide_layouts[5]  # Try standard blank layout
        except IndexError:
            slide_layout = prs.slide_layouts[-1]  # Use last available layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = section.get('title', 'Analysis')
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        
        # Add content on left side
        content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(5))
        content_frame = content_shape.text_frame
        
        for i, point in enumerate(section.get('content', [])):
            if i == 0:
                content_frame.text = f"• {point}"
            else:
                p = content_frame.add_paragraph()
                p.text = f"• {point}"
            
        # Could add chart or image on right side
        self._apply_colors(slide, colors)
    
    def _fill_remaining_slides(self, prs: Presentation, structure: Dict, current: int, target: int):
        """Fill remaining slides to reach target count"""
        
        # Additional sections for PE/IB decks
        additional_sections = [
            {"title": "Financial Projections", "content": ["5-year financial model", "Revenue assumptions", "Cost structure", "Capital requirements"]},
            {"title": "Comparable Analysis", "content": ["Trading comparables", "Transaction comparables", "Valuation benchmarks"]},
            {"title": "Due Diligence Findings", "content": ["Financial due diligence", "Legal due diligence", "Operational review", "ESG considerations"]},
            {"title": "Synergy Analysis", "content": ["Revenue synergies", "Cost synergies", "Timeline to realization"]},
            {"title": "Integration Plan", "content": ["100-day plan", "Key milestones", "Integration team", "Risk mitigation"]},
            {"title": "Regulatory Considerations", "content": ["Regulatory approvals", "Compliance requirements", "Timeline"]},
            {"title": "Financing Structure", "content": ["Sources of funds", "Uses of funds", "Pro forma capitalization"]},
            {"title": "Management Incentives", "content": ["Equity participation", "Performance targets", "Retention strategy"]},
            {"title": "ESG Considerations", "content": ["Environmental impact", "Social responsibility", "Governance structure"]},
            {"title": "Appendix", "content": ["Detailed financials", "Supporting documentation", "Contact information"]}
        ]
        
        slides_needed = target - current
        for i in range(min(slides_needed, len(additional_sections))):
            section = additional_sections[i]
            section['slide_type'] = 'content'
            self._create_content_slide(prs, section, structure['color_scheme'])
    
    def _apply_colors(self, slide, colors: Dict):
        """Apply color scheme to slide elements"""
        # This would apply the color scheme to various elements
        # For now, it's a placeholder for future enhancement
        pass