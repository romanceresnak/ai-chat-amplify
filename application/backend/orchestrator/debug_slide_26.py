#!/usr/bin/env python3
"""
Debug script to investigate Slide 26 generation issue
"""

import boto3
import io
import logging
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(message)s')
logger = logging.getLogger(__name__)

# Initialize S3
s3 = boto3.client('s3')

def analyze_shape(shape, indent=""):
    """Analyze and print details about a shape"""
    info = []
    info.append(f"{indent}Shape: {shape.name}")
    info.append(f"{indent}  Type: {shape.shape_type}")
    
    if hasattr(shape, 'text'):
        text = shape.text.strip()
        if text:
            info.append(f"{indent}  Text: '{text[:50]}{'...' if len(text) > 50 else ''}'")
    
    if hasattr(shape, 'has_chart') and shape.has_chart:
        info.append(f"{indent}  Has Chart: Yes")
        info.append(f"{indent}  Chart Type: {shape.chart.chart_type}")
    
    if hasattr(shape, 'has_table') and shape.has_table:
        info.append(f"{indent}  Has Table: Yes")
        info.append(f"{indent}  Table dimensions: {shape.table.rows} x {shape.table.columns}")
    
    # Check if it's a group shape
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        info.append(f"{indent}  Group with {len(shape.shapes)} shapes")
        for sub_shape in shape.shapes:
            sub_info = analyze_shape(sub_shape, indent + "    ")
            info.extend(sub_info)
    
    return info

def find_slide_26_in_template():
    """Find and analyze slide 26 in the template"""
    logger.info("=" * 80)
    logger.info("ANALYZING TEMPLATE FOR SLIDE 26")
    logger.info("=" * 80)
    
    # Load template from S3
    try:
        response = s3.get_object(
            Bucket='scribbe-ai-dev-documents',
            Key='PUBLIC IP South Plains (1).pptx'
        )
        template_bytes = response['Body'].read()
        logger.info(f"Loaded template: {len(template_bytes)} bytes")
    except Exception as e:
        logger.error(f"Failed to load template: {e}")
        return
    
    # Load presentation
    prs = Presentation(io.BytesIO(template_bytes))
    logger.info(f"Total slides in template: {len(prs.slides)}")
    
    # Look for slide with "Loan Portfolio" title
    loan_portfolio_slides = []
    
    for idx, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_text.append(shape.text)
        
        full_text = ' '.join(slide_text).lower()
        
        if 'loan portfolio' in full_text:
            logger.info(f"\nFound 'Loan Portfolio' at slide position {idx + 1}")
            
            # Check if it has the specific content markers
            if '$1,936' in full_text:
                logger.info(f"  -> Contains $1,936 (Slide 23/26 content)")
                loan_portfolio_slides.append({
                    'position': idx + 1,
                    'type': 'bar_chart',
                    'slide': slide
                })
            elif 'commercial real estate' in full_text and '%' in full_text:
                logger.info(f"  -> Contains portfolio percentages (Slide 24 content)")
                loan_portfolio_slides.append({
                    'position': idx + 1,
                    'type': 'donut_chart',
                    'slide': slide
                })
    
    # Analyze slides that match Slide 23/26 pattern
    logger.info("\n" + "=" * 80)
    logger.info("DETAILED ANALYSIS OF SLIDE 23/26 CONTENT")
    logger.info("=" * 80)
    
    for slide_info in loan_portfolio_slides:
        if slide_info['type'] == 'bar_chart':
            slide = slide_info['slide']
            position = slide_info['position']
            
            logger.info(f"\nAnalyzing slide at position {position}:")
            logger.info(f"Number of shapes: {len(slide.shapes)}")
            
            # Analyze each shape
            for i, shape in enumerate(slide.shapes):
                logger.info(f"\nShape {i}:")
                shape_info = analyze_shape(shape, "  ")
                for line in shape_info:
                    logger.info(line)
            
            # Check specifically for chart
            has_chart = False
            for shape in slide.shapes:
                if hasattr(shape, 'has_chart') and shape.has_chart:
                    has_chart = True
                    logger.info("\n*** FOUND CHART IN SLIDE ***")
                    logger.info(f"Chart type: {shape.chart.chart_type}")
                    logger.info(f"Shape name: {shape.name}")
                    logger.info(f"Shape position: Left={shape.left}, Top={shape.top}")
                    logger.info(f"Shape size: Width={shape.width}, Height={shape.height}")
            
            if not has_chart:
                logger.warning("\n*** NO CHART FOUND IN SLIDE ***")
    
    # Test the SmartTemplateGenerator update process
    logger.info("\n" + "=" * 80)
    logger.info("TESTING SMARTTEMPLATEGENERATOR UPDATE PROCESS")
    logger.info("=" * 80)
    
    from smart_template_generator import SmartTemplateGenerator
    
    generator = SmartTemplateGenerator()
    
    # Test instruction
    test_instruction = """Generate Slide 26 with updated values: $2,000M for 2Q'19, $2,100M for 3Q'19, 
    $2,200M for 4Q'19, $2,300M for 1Q'20, $2,500M for 2Q'20. 
    Also update yield percentages (6.0%, 6.1%, 5.9%, 5.8%, 5.5%)"""
    
    try:
        # Parse the instruction
        slide_info = generator._parse_instructions(test_instruction)
        logger.info(f"\nParsed instruction: {slide_info}")
        
        # Find the slide to update
        target_slide = None
        for idx, slide in enumerate(prs.slides):
            slide_text = generator._get_slide_text(slide).lower()
            if 'loan portfolio' in slide_text and '$1,936' in slide_text:
                target_slide = slide
                logger.info(f"\nFound target slide at position {idx + 1}")
                break
        
        if target_slide:
            logger.info("\nBefore update - Shape texts:")
            shape_count = 0
            for shape in target_slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    logger.info(f"  Shape {shape_count}: '{shape.text}'")
                shape_count += 1
                if shape_count >= 15:  # Only show first 15
                    break
            
            # Update the slide
            generator._update_slide_23(target_slide, slide_info)
            
            logger.info("\nAfter update - Shape texts:")
            shape_count = 0
            for shape in target_slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    logger.info(f"  Shape {shape_count}: '{shape.text}'")
                shape_count += 1
                if shape_count >= 15:  # Only show first 15
                    break
            
            # Check if chart still exists
            has_chart_after = False
            for shape in target_slide.shapes:
                if hasattr(shape, 'has_chart') and shape.has_chart:
                    has_chart_after = True
                    logger.info("\n*** CHART STILL EXISTS AFTER UPDATE ***")
                    break
            
            if not has_chart_after:
                logger.warning("\n*** CHART MISSING AFTER UPDATE ***")
    
    except Exception as e:
        logger.error(f"Error during update test: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    find_slide_26_in_template()