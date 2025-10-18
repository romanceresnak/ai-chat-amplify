"""
AI-Powered Presentation Generator - Template Based Version
Uses existing PPTX template to avoid corruption issues
"""

import json
import boto3
from typing import Dict, List, Any
import logging
import io
import re
import os

logger = logging.getLogger()
logger.setLevel(logging.INFO)

# Initialize S3 client
s3 = boto3.client('s3')

class AIPresentationGenerator:
    def __init__(self):
        self.bedrock_runtime = boto3.client('bedrock-runtime', region_name='us-east-1')
        self.model_id = 'eu.anthropic.claude-3-5-sonnet-20240620-v1:0'
        self.documents_bucket = os.environ.get('DOCUMENTS_BUCKET', 'scribbe-ai-dev-documents')
        self.template_key = 'PUBLIC IP South Plains (1).pptx'  # Use existing template
    
    def generate_presentation(self, instructions: str) -> bytes:
        """Generate a presentation by modifying existing template"""
        
        try:
            # Import python-pptx
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.chart.data import ChartData
            from pptx.enum.chart import XL_CHART_TYPE
            
            logger.info("Successfully imported python-pptx")
            
        except ImportError as e:
            logger.error(f"Failed to import python-pptx: {e}")
            # Fallback to creating minimal XML-based presentation
            return self._create_minimal_presentation()
        
        try:
            # Download template from S3
            logger.info(f"Downloading template: {self.template_key}")
            template_response = s3.get_object(Bucket=self.documents_bucket, Key=self.template_key)
            template_content = template_response['Body'].read()
            
            # Open template presentation
            template_stream = io.BytesIO(template_content)
            prs = Presentation(template_stream)
            
            logger.info(f"Loaded template with {len(prs.slides)} slides")
            
            # Keep first slide as title
            if len(prs.slides) > 0:
                title_slide = prs.slides[0]
                if title_slide.shapes.title:
                    title_slide.shapes.title.text = "Loan Portfolio Analysis"
                
                # Update subtitle if exists
                for shape in title_slide.shapes:
                    if shape.has_text_frame and shape != title_slide.shapes.title:
                        shape.text = "South Plains Financial, Inc.\nSeptember 2024"
                        break
            
            # Remove all slides except first one
            while len(prs.slides) > 1:
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                xml_slides.remove(slides[-1])
            
            # Add new content slides
            self._add_portfolio_slides(prs, Inches, Pt, RGBColor, ChartData, XL_CHART_TYPE)
            
            # Save to bytes using getvalue() 
            output = io.BytesIO()
            prs.save(output)
            pptx_content = output.getvalue()
            
            logger.info(f"Generated presentation with {len(prs.slides)} slides, size: {len(pptx_content)} bytes")
            
            return pptx_content
            
        except Exception as e:
            logger.error(f"Error generating presentation: {e}")
            # Fallback to minimal presentation
            return self._create_minimal_presentation()
    
    def _add_portfolio_slides(self, prs, Inches, Pt, RGBColor, ChartData, XL_CHART_TYPE):
        """Add loan portfolio slides to presentation"""
        
        # Slide 2: Portfolio Overview
        slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = "Loan Portfolio Overview"
        
        # Add content
        try:
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.has_text_frame and shape != slide.shapes.title:
                    content_placeholder = shape
                    break
            
            if content_placeholder:
                tf = content_placeholder.text_frame
                tf.text = "Net Loans: $2.3 Billion"
                p = tf.add_paragraph()
                p.text = "Commercial Real Estate: 28%"
                p = tf.add_paragraph()
                p.text = "Commercial General: 27%"
                p = tf.add_paragraph()
                p.text = "Residential: 15%"
                p = tf.add_paragraph()
                p.text = "Auto Loans: 16%"
        except:
            pass
        
        # Slide 3: Key Highlights
        slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = "Key Highlights"
        
        try:
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.has_text_frame and shape != slide.shapes.title:
                    content_placeholder = shape
                    break
            
            if content_placeholder:
                tf = content_placeholder.text_frame
                tf.text = "Strong asset quality metrics"
                p = tf.add_paragraph()
                p.text = "Diversified loan portfolio"
                p = tf.add_paragraph()
                p.text = "Conservative underwriting standards"
                p = tf.add_paragraph()
                p.text = "Experienced management team"
        except:
            pass
    
    def _create_minimal_presentation(self) -> bytes:
        """Create minimal presentation without python-pptx"""
        import zipfile
        import xml.etree.ElementTree as ET
        
        # Create minimal PPTX structure
        output = io.BytesIO()
        
        with zipfile.ZipFile(output, 'w', zipfile.ZIP_DEFLATED) as zipf:
            # Add content types
            content_types = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
    <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
    <Default Extension="xml" ContentType="application/xml"/>
    <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
    <Override PartName="/ppt/slides/slide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
    <Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>
    <Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>
</Types>'''
            zipf.writestr('[Content_Types].xml', content_types)
            
            # Add _rels/.rels
            rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
</Relationships>'''
            zipf.writestr('_rels/.rels', rels)
            
            # Add ppt/presentation.xml
            presentation = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
    <p:sldMasterIdLst>
        <p:sldMasterId id="2147483648" r:id="rId1"/>
    </p:sldMasterIdLst>
    <p:sldIdLst>
        <p:sldId id="256" r:id="rId2"/>
    </p:sldIdLst>
    <p:sldSz cx="9144000" cy="6858000"/>
    <p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>'''
            zipf.writestr('ppt/presentation.xml', presentation)
            
            # Add ppt/_rels/presentation.xml.rels
            ppt_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>
    <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/>
</Relationships>'''
            zipf.writestr('ppt/_rels/presentation.xml.rels', ppt_rels)
            
            # Add minimal slide
            slide1 = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
    <p:cSld>
        <p:spTree>
            <p:nvGrpSpPr>
                <p:cNvPr id="1" name=""/>
                <p:cNvGrpSpPr/>
                <p:nvPr/>
            </p:nvGrpSpPr>
            <p:grpSpPr>
                <a:xfrm>
                    <a:off x="0" y="0"/>
                    <a:ext cx="0" cy="0"/>
                    <a:chOff x="0" y="0"/>
                    <a:chExt cx="0" cy="0"/>
                </a:xfrm>
            </p:grpSpPr>
            <p:sp>
                <p:nvSpPr>
                    <p:cNvPr id="2" name="Title"/>
                    <p:cNvSpPr>
                        <a:spLocks noGrp="1"/>
                    </p:cNvSpPr>
                    <p:nvPr>
                        <p:ph type="title"/>
                    </p:nvPr>
                </p:nvSpPr>
                <p:spPr/>
                <p:txBody>
                    <a:bodyPr/>
                    <a:lstStyle/>
                    <a:p>
                        <a:r>
                            <a:rPr lang="en-US"/>
                            <a:t>Loan Portfolio Analysis</a:t>
                        </a:r>
                    </a:p>
                </p:txBody>
            </p:sp>
        </p:spTree>
    </p:cSld>
    <p:clrMapOvr>
        <a:masterClrMapping/>
    </p:clrMapOvr>
</p:sld>'''
            zipf.writestr('ppt/slides/slide1.xml', slide1)
            
            # Add slide relationship
            slide_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
</Relationships>'''
            zipf.writestr('ppt/slides/_rels/slide1.xml.rels', slide_rels)
            
            # Add minimal slide layout
            slide_layout = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="title" preserve="1">
    <p:cSld name="Title Slide">
        <p:spTree>
            <p:nvGrpSpPr>
                <p:cNvPr id="1" name=""/>
                <p:cNvGrpSpPr/>
                <p:nvPr/>
            </p:nvGrpSpPr>
            <p:grpSpPr/>
        </p:spTree>
    </p:cSld>
    <p:clrMapOvr>
        <a:masterClrMapping/>
    </p:clrMapOvr>
</p:sldLayout>'''
            zipf.writestr('ppt/slideLayouts/slideLayout1.xml', slide_layout)
            
            # Add layout relationship
            layout_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/>
</Relationships>'''
            zipf.writestr('ppt/slideLayouts/_rels/slideLayout1.xml.rels', layout_rels)
            
            # Add minimal slide master
            slide_master = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
    <p:cSld>
        <p:spTree>
            <p:nvGrpSpPr>
                <p:cNvPr id="1" name=""/>
                <p:cNvGrpSpPr/>
                <p:nvPr/>
            </p:nvGrpSpPr>
            <p:grpSpPr/>
        </p:spTree>
    </p:cSld>
    <p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/>
    <p:sldLayoutIdLst>
        <p:sldLayoutId id="2147483649" r:id="rId1"/>
    </p:sldLayoutIdLst>
</p:sldMaster>'''
            zipf.writestr('ppt/slideMasters/slideMaster1.xml', slide_master)
            
            # Add master relationship
            master_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
</Relationships>'''
            zipf.writestr('ppt/slideMasters/_rels/slideMaster1.xml.rels', master_rels)
        
        output.seek(0)
        return output.getvalue()

    def analyze_presentation_request(self, instructions: str) -> Dict[str, Any]:
        """Simple analysis - returns loan portfolio structure"""
        return self._get_loan_portfolio_structure()
    
    def _get_loan_portfolio_structure(self) -> Dict[str, Any]:
        """Loan portfolio presentation structure"""
        return {
            "presentation_type": "loan_portfolio",
            "title": "Loan Portfolio Analysis",
            "sections": []
        }
    
    def _get_default_structure(self, instructions: str) -> Dict[str, Any]:
        """Default structure"""
        return self._get_loan_portfolio_structure()