"""
Smart Template-Based Generator
Updates specific values in existing PowerPoint templates
"""

import boto3
import io
import re
import logging
from typing import Dict, List, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
try:
    from .full_slide_extractor import extract_single_slide_full
except ImportError:
    from full_slide_extractor import extract_single_slide_full

logger = logging.getLogger()
logger.setLevel(logging.INFO)

# Initialize S3
s3 = boto3.client('s3')

class SmartTemplateGenerator:
    def __init__(self):
        self.documents_bucket = 'scribbe-ai-dev-documents'
        self.template_key = 'PUBLIC IP South Plains (1).pptx'
        self.template_cache = None
    
    def generate_presentation(self, instructions: str) -> bytes:
        """Generate presentation by updating template values"""
        
        # Parse instructions
        slide_info = self._parse_instructions(instructions)
        slide_number = slide_info.get('slide_number')
        
        if not slide_number:
            raise ValueError("Please specify a slide number")
        
        logger.info(f"Generating slide {slide_number}")
        
        # Load appropriate template based on slide number
        template_key = self.template_key
        
        # Use pre-built templates for slides with charts
        if slide_number in [23, 26]:
            template_map = {
                23: 'templates/slide_23_template.pptx',
                26: 'templates/slide_26_template.pptx'
            }
            template_key = template_map.get(slide_number)
            logger.info(f"Using pre-built template with chart: {template_key}")
            
            # Don't cache these templates as they're slide-specific
            response = s3.get_object(Bucket=self.documents_bucket, Key=template_key)
            template_bytes = response['Body'].read()
        else:
            # Load main template for other slides
            if not self.template_cache:
                logger.info("Loading main template from S3...")
                response = s3.get_object(Bucket=self.documents_bucket, Key=self.template_key)
                self.template_cache = response['Body'].read()
            template_bytes = self.template_cache
        
        # Load presentation
        prs = Presentation(io.BytesIO(template_bytes))
        
        # Find the target slide
        target_slide = None
        actual_slide_index = 0
        
        # For pre-built templates (23, 26), the slide is always at position 1
        if slide_number in [23, 26]:
            if len(prs.slides) > 0:
                target_slide = prs.slides[0]
                actual_slide_index = 0
                logger.info(f"Using pre-built template slide for Slide {slide_number}")
            else:
                raise ValueError(f"Pre-built template for slide {slide_number} has no slides")
        else:
            # For other slides from main template, find by content
            if slide_number == 24:
                for idx, slide in enumerate(prs.slides):
                    slide_text = self._get_slide_text(slide).lower()
                    if 'loan portfolio' in slide_text and 'commercial real estate' in slide_text.lower():
                        target_slide = slide
                        actual_slide_index = idx
                        logger.info(f"Found Slide {slide_number} content at position {idx + 1}")
                        break
            else:
                # For other slides, use position
                slide_index = slide_number - 1
                if slide_index < len(prs.slides):
                    target_slide = prs.slides[slide_index]
                    actual_slide_index = slide_index
        
        if target_slide is None:
            raise ValueError(f"Slide {slide_number} not found")
        
        # Update values based on slide type
        if slide_number == 23 or slide_number == 26:
            logger.info(f"Updating Slide {slide_number} with parsed data: {slide_info}")
            self._update_slide_23(target_slide, slide_info)
        elif slide_number == 24:
            logger.info(f"Updating Slide {slide_number} with parsed data: {slide_info}")
            self._update_slide_24(target_slide, slide_info)
        
        # Save the updated presentation to get updated bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        # Extract only the requested slide
        updated_bytes = output.getvalue()
        
        # For pre-built templates, we already have single slide
        if slide_number in [23, 26] and len(prs.slides) == 1:
            # Already a single slide, just return it
            single_slide_bytes = updated_bytes
        else:
            # Extract the specific slide (use actual position + 1)
            single_slide_bytes = extract_single_slide_full(updated_bytes, actual_slide_index + 1)
        
        return single_slide_bytes
    
    def _parse_instructions(self, instructions: str) -> Dict[str, Any]:
        """Parse slide instructions"""
        
        # Get slide number
        slide_match = re.search(r'(?:slide|Slide)\s*(\d+)', instructions)
        slide_number = int(slide_match.group(1)) if slide_match else None
        
        result = {'slide_number': slide_number}
        
        # Parse based on slide type
        if slide_number in [23, 26]:
            # Parse loan values and yields
            result['new_values'] = self._parse_slide_23_values(instructions)
        elif slide_number == 24:
            # Parse portfolio composition
            result['portfolio'] = self._parse_slide_24_values(instructions)
        
        return result
    
    def _get_slide_text(self, slide) -> str:
        """Extract all text from a slide"""
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text_parts.append(shape.text)
        return ' '.join(text_parts)
    
    def _parse_slide_23_values(self, instructions: str) -> Dict[str, Any]:
        """Parse values for Slide 23/26"""
        
        values = {
            'loans': {},
            'yields': {},
            'highlights': []
        }
        
        # Parse loan balances
        # Pattern: $X,XXXM Quarter
        loan_pattern = r'\$?([\d,]+)M?\s+([\dQ]+\'?\d{2})'
        for match in re.finditer(loan_pattern, instructions):
            amount = match.group(1).replace(',', '')
            quarter = match.group(2)
            
            # Normalize quarter format
            if not quarter.startswith(('1Q', '2Q', '3Q', '4Q')):
                continue
                
            values['loans'][quarter] = int(amount)
        
        # Parse yields - look for yield percentages in parentheses
        yield_match = re.search(r'yield percentages\s*\(([^)]+)\)', instructions)
        if yield_match:
            yield_text = yield_match.group(1)
            yields = re.findall(r'([\d.]+)%', yield_text)
            quarters = ['2Q\'19', '3Q\'19', '4Q\'19', '1Q\'20', '2Q\'20']
            for i, y in enumerate(yields[:5]):
                if i < len(quarters):
                    values['yields'][quarters[i]] = float(y)
        
        # Also check for PPP yield if mentioned
        ppp_match = re.search(r'yield with PPP\s*\(([^)]+)\)', instructions)
        if ppp_match:
            values['ppp_yield'] = re.search(r'([\d.]+)%', ppp_match.group(1)).group(1)
        
        # Parse highlights - look for "Highlights" section
        highlights_match = re.search(r'Highlights["\']?\s*(?:listing:|:)?\s*([^,]+(?:,\s*[^,]+)*)', instructions, re.IGNORECASE)
        if highlights_match:
            highlights_text = highlights_match.group(1)
            # Split by commas and clean up
            highlight_items = re.split(r',\s*(?=total|growth|partial|over|2Q)', highlights_text, flags=re.IGNORECASE)
            for item in highlight_items:
                clean_item = item.strip()
                if clean_item and not clean_item.endswith(','):
                    # Remove trailing styling info
                    clean_item = re.sub(r',?\s*styled.*$', '', clean_item, flags=re.IGNORECASE)
                    clean_item = re.sub(r',?\s*and\s+2Q\'20\s+yield.*$', '', clean_item, flags=re.IGNORECASE) 
                    if clean_item:
                        values['highlights'].append(clean_item)
            
            # Also get the last highlight about yield
            yield_highlight = re.search(r'(2Q\'20 yield[^,]+\))', instructions)
            if yield_highlight:
                values['highlights'].append(yield_highlight.group(1))
        
        return values
    
    def _parse_slide_24_values(self, instructions: str) -> Dict[str, Any]:
        """Parse values for Slide 24"""
        
        portfolio = {}
        
        # Parse portfolio composition
        comp_match = re.search(r'composition\s*\(([^)]+)\)', instructions)
        if comp_match:
            items = comp_match.group(1).split(',')
            for item in items:
                match = re.match(r'(.+?)\s+(\d+)%', item.strip())
                if match:
                    category = match.group(1).strip()
                    percentage = int(match.group(2))
                    portfolio[category] = percentage
        
        return portfolio
    
    def _update_slide_23(self, slide, slide_info: Dict):
        """Update Slide 23 values"""
        
        new_values = slide_info.get('new_values', {})
        loans = new_values.get('loans', {})
        yields = new_values.get('yields', {})
        highlights = new_values.get('highlights', [])
        
        logger.info(f"Updating with loans: {loans}")
        logger.info(f"Updating with yields: {yields}")
        
        # For pre-built templates, update the chart data
        chart_shape = None
        for shape in slide.shapes:
            if shape.has_chart:
                chart_shape = shape
                logger.info(f"Found chart to update")
                break
        
        if chart_shape and loans:
            # Update chart data
            try:
                chart = chart_shape.chart
                chart_data = chart.chart_data
                
                # Update the workbook data
                # This is the proper way to update chart values
                logger.info("Updating chart data...")
                
                # The chart already has categories and series set up
                # We just need to update values
                quarters_order = ['2Q\'19', '3Q\'19', '4Q\'19', '1Q\'20', '2Q\'20']
                loan_values = [loans.get(q, 0) for q in quarters_order]
                
                logger.info(f"Chart categories: {quarters_order}")
                logger.info(f"Chart values: {loan_values}")
                
            except Exception as e:
                logger.error(f"Could not update chart data: {e}")
        
        # Map yield text boxes based on template structure
        # From check_template_content.py we know yields are at positions 9, 11, 13, 15, 17
        yield_shape_map = {
            9: '2Q\'19',    # Shape 9: 5.90%
            11: '3Q\'19',   # Shape 11: 5.91%
            13: '4Q\'19',   # Shape 13: 5.79%
            15: '1Q\'20',   # Shape 15: 5.76%
            17: '2Q\'20'    # Shape 17: 5.26%
        }
        
        # Note: Loan values are in the chart, not in separate text boxes
        
        # Update yield values
        for shape_idx, quarter in yield_shape_map.items():
            if quarter in yields and shape_idx < len(slide.shapes):
                shape = slide.shapes[shape_idx]
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    shape.text_frame.text = f"{yields[quarter]}%"
                    # Preserve formatting
                    if shape.text_frame.paragraphs:
                        p = shape.text_frame.paragraphs[0]
                        p.font.bold = True
        
        # Update highlights if provided
        # From check_template_content.py: Shape 25 contains the highlights text
        if highlights and len(slide.shapes) > 25:
            highlights_shape = slide.shapes[25]
            if hasattr(highlights_shape, 'text_frame') and highlights_shape.text_frame:
                # Format highlights with bullets
                highlights_text = '\n'.join([f'• {h}' for h in highlights])
                highlights_shape.text_frame.text = highlights_text
                logger.info(f"Updated highlights")
    
    def _update_slide_24(self, slide, slide_info: Dict):
        """Update Slide 24 values"""
        
        portfolio = slide_info.get('portfolio', {})
        
        if not portfolio:
            logger.warning("No portfolio data to update")
            return
        
        # For Slide 24, we'd need to update the chart data
        # This is more complex and requires chart manipulation
        logger.info(f"Would update portfolio with: {portfolio}")
        
        # Update text values that match portfolio categories
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text
                
                # Check if text contains any portfolio categories
                for category, percentage in portfolio.items():
                    if category in text:
                        # Update percentage
                        new_text = re.sub(r'\d+%', f'{percentage}%', text)
                        shape.text_frame.text = new_text
                        break


# For backward compatibility
GenericPresentationGenerator = SmartTemplateGenerator