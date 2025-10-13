#!/usr/bin/env python3
"""
Test script to verify python-pptx imports in Lambda environment
"""
import os
import sys
import logging

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def test_lxml_import():
    """Test lxml import specifically"""
    try:
        logger.info("Testing lxml import...")
        import lxml
        logger.info(f"lxml imported successfully, version: {lxml.__version__}")
        
        # Test etree specifically
        from lxml import etree
        logger.info("lxml.etree imported successfully")
        
        # Test basic etree functionality
        root = etree.Element("test")
        logger.info("lxml.etree basic functionality works")
        
        return True
    except Exception as e:
        logger.error(f"lxml import failed: {e}")
        return False

def test_pptx_import():
    """Test python-pptx import"""
    try:
        logger.info("Testing python-pptx import...")
        
        # Add layer path if needed
        layer_path = "/opt/python/lib/python3.11/site-packages"
        if os.path.exists(layer_path) and layer_path not in sys.path:
            sys.path.insert(0, layer_path)
            logger.info(f"Added {layer_path} to sys.path")
        
        # Test core imports
        from pptx import Presentation
        logger.info("pptx.Presentation imported successfully")
        
        from pptx.util import Inches, Pt
        logger.info("pptx.util imports successful")
        
        from pptx.dml.color import RGBColor
        logger.info("pptx.dml.color.RGBColor imported successfully")
        
        from pptx.enum.text import PP_ALIGN
        logger.info("pptx.enum.text.PP_ALIGN imported successfully")
        
        from pptx.chart.data import ChartData
        logger.info("pptx.chart.data.ChartData imported successfully")
        
        from pptx.enum.chart import XL_CHART_TYPE
        logger.info("pptx.enum.chart.XL_CHART_TYPE imported successfully")
        
        # Test basic functionality
        prs = Presentation()
        logger.info("Created presentation object successfully")
        
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        logger.info("Added slide successfully")
        
        return True
        
    except Exception as e:
        logger.error(f"python-pptx import failed: {e}")
        logger.error(f"Python path: {sys.path}")
        return False

def lambda_handler(event, context):
    """Lambda entry point for testing"""
    logger.info("Starting python-pptx import test...")
    
    # Test lxml first
    lxml_success = test_lxml_import()
    
    # Test python-pptx
    pptx_success = test_pptx_import()
    
    result = {
        "lxml_import": lxml_success,
        "pptx_import": pptx_success,
        "overall_success": lxml_success and pptx_success
    }
    
    logger.info(f"Test results: {result}")
    
    return {
        "statusCode": 200,
        "body": result
    }

if __name__ == "__main__":
    # For local testing
    lambda_handler({}, {})