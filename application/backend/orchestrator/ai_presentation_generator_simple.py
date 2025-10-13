"""
AI-Powered Presentation Generator for Financial Services - Simplified Version
Generates professional presentations using Bedrock AI and python-pptx with standard layouts
"""

import json
import boto3
from typing import Dict, List, Any, Optional
import logging
import io
import re

logger = logging.getLogger()

# Global variables for python-pptx availability
PPTX_AVAILABLE = False
Presentation = None
Inches = None
Pt = None
RGBColor = None
PP_ALIGN = None
ChartData = None
XL_CHART_TYPE = None
XL_LEGEND_POSITION = None

class AIPresentationGenerator:
    def __init__(self):
        self.bedrock_runtime = boto3.client('bedrock-runtime', region_name='eu-west-1')
        self.model_id = 'eu.anthropic.claude-3-5-sonnet-20240620-v1:0'
        self._initialize_pptx()
    
    def _initialize_pptx(self):
        """Initialize python-pptx imports on demand"""
        global PPTX_AVAILABLE, Presentation, Inches, Pt, RGBColor, PP_ALIGN, ChartData, XL_CHART_TYPE, XL_LEGEND_POSITION
        
        if PPTX_AVAILABLE:
            return  # Already initialized
        
        try:
            import sys
            import os
            
            # Debug info
            logger.info(f"Python path: {sys.path}")
            layer_path = "/opt/python/lib/python3.11/site-packages"
            
            if os.path.exists(layer_path):
                logger.info(f"Layer directory exists: {layer_path}")
                if layer_path not in sys.path:
                    sys.path.insert(0, layer_path)
                    logger.info(f"Added {layer_path} to sys.path")
            
            # Import python-pptx components
            from pptx import Presentation as _Presentation
            from pptx.util import Inches as _Inches, Pt as _Pt
            from pptx.dml.color import RGBColor as _RGBColor
            from pptx.enum.text import PP_ALIGN as _PP_ALIGN
            from pptx.chart.data import ChartData as _ChartData
            from pptx.enum.chart import XL_CHART_TYPE as _XL_CHART_TYPE, XL_LEGEND_POSITION as _XL_LEGEND_POSITION
            
            # Set global variables
            Presentation = _Presentation
            Inches = _Inches
            Pt = _Pt
            RGBColor = _RGBColor
            PP_ALIGN = _PP_ALIGN
            ChartData = _ChartData
            XL_CHART_TYPE = _XL_CHART_TYPE
            XL_LEGEND_POSITION = _XL_LEGEND_POSITION
            PPTX_AVAILABLE = True
            
            logger.info("python-pptx successfully imported in AI generator")
            
        except ImportError as e:
            logger.error(f"Failed to import python-pptx: {e}")
            PPTX_AVAILABLE = False
        
    def analyze_presentation_request(self, instructions: str) -> Dict[str, Any]:
        """Use AI to analyze and structure the presentation request"""
        
        prompt = f"""Analyze this presentation request and provide a structured plan.
        
Request: {instructions}

Return a JSON structure with:
1. presentation_type: "financial_pe", "financial_ib", "loan_portfolio", "general"
2. title: Main presentation title
3. num_slides: Number of slides to create (default 40 for PE/IB decks)
4. sections: Array of sections, each containing:
   - title: Section title
   - slide_type: "title", "content", "chart", "table", "mixed"
   - content: Key points or data
   - chart_data: If applicable, data for charts
5. design_theme: "corporate", "financial", "modern"
6. color_scheme: Primary colors to use

For financial presentations, include standard sections like:
- Executive Summary
- Company Overview
- Financial Performance
- Market Analysis
- Investment Thesis
- Risk Factors
- Appendices

Be specific and detailed."""

        # Use the correct format for Claude 3.5 Sonnet
        request_body = {
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 2000,
            "temperature": 0.7,
            "messages": [
                {
                    "role": "user",
                    "content": f"{prompt}\n\nReturn only valid JSON in your response."
                }
            ]
        }
        
        response = self.bedrock_runtime.invoke_model(
            modelId=self.model_id,
            body=json.dumps(request_body),
            contentType='application/json'
        )
        
        response_body = json.loads(response['body'].read())
        response_text = response_body.get('content', [{}])[0].get('text', '')
        
        # Extract JSON from response
        json_match = re.search(r'```json\n(.*?)\n```', response_text, re.DOTALL)
        if json_match:
            return json.loads(json_match.group(1))
        else:
            # Fallback structure
            return self._get_default_structure(instructions)
    
    def _get_default_structure(self, instructions: str) -> Dict[str, Any]:
        """Provide default structure based on keywords"""
        
        instructions_lower = instructions.lower()
        
        if "private equity" in instructions_lower or "investment committee" in instructions_lower:
            return self._get_pe_deck_structure()
        elif "debt issuance" in instructions_lower:
            return self._get_debt_issuance_structure()
        elif "loan portfolio" in instructions_lower:
            return self._get_loan_portfolio_structure()
        else:
            return self._get_general_financial_structure()
    
    def _get_pe_deck_structure(self) -> Dict[str, Any]:
        """Standard PE investment committee deck structure"""
        return {
            "presentation_type": "financial_pe",
            "title": "Private Equity Investment Committee Presentation",
            "num_slides": 40,
            "design_theme": "financial",
            "color_scheme": {"primary": "002855", "accent": "CC0000"},
            "sections": [
                {
                    "title": "Title Slide",
                    "slide_type": "title",
                    "content": ["Investment Committee Meeting", "Confidential"]
                },
                {
                    "title": "Executive Summary",
                    "slide_type": "content",
                    "content": [
                        "Investment opportunity overview",
                        "Key investment highlights",
                        "Expected returns and timeline",
                        "Recommendation"
                    ]
                },
                {
                    "title": "Investment Thesis",
                    "slide_type": "mixed",
                    "content": [
                        "Strategic rationale",
                        "Value creation opportunities",
                        "Competitive advantages",
                        "Market positioning"
                    ]
                },
                {
                    "title": "Company Overview",
                    "slide_type": "content",
                    "content": [
                        "Business description",
                        "Products and services",
                        "Management team",
                        "Corporate structure"
                    ]
                },
                {
                    "title": "Financial Performance",
                    "slide_type": "chart",
                    "content": ["Historical financials", "Revenue growth", "EBITDA margins"],
                    "chart_data": {
                        "type": "column",
                        "categories": ["2020", "2021", "2022", "2023"],
                        "series": [
                            {"name": "Revenue", "values": [100, 120, 145, 170]},
                            {"name": "EBITDA", "values": [20, 28, 35, 42]}
                        ]
                    }
                }
            ]
        }
    
    def _get_debt_issuance_structure(self) -> Dict[str, Any]:
        """Debt issuance presentation structure"""
        return {
            "presentation_type": "financial_ib",
            "title": "Debt Issuance Presentation",
            "num_slides": 35,
            "design_theme": "corporate",
            "color_scheme": {"primary": "003366", "accent": "FF6600"},
            "sections": [
                {
                    "title": "Title Slide",
                    "slide_type": "title",
                    "content": ["Debt Issuance Opportunity", "Institutional Investor Presentation"]
                },
                {
                    "title": "Transaction Overview",
                    "slide_type": "content",
                    "content": [
                        "Issuer overview",
                        "Transaction size and structure",
                        "Use of proceeds",
                        "Key terms"
                    ]
                }
            ]
        }
    
    def _get_loan_portfolio_structure(self) -> Dict[str, Any]:
        """Loan portfolio presentation structure with professional design"""
        return {
            "presentation_type": "loan_portfolio",
            "title": "Loan Portfolio Analysis",
            "num_slides": 15,
            "design_theme": "modern",
            "color_scheme": {"primary": "1F497D", "accent": "C0504D"},
            "sections": [
                {
                    "title": "South Plains Financial, Inc.",
                    "slide_type": "title",
                    "content": ["September 2024", "Loan Portfolio Analysis", "Investor Presentation"]
                },
                {
                    "title": "Loan Portfolio",
                    "slide_type": "chart",
                    "content": ["Portfolio composition overview", "Diversified loan book across multiple segments"],
                    "chart_data": {
                        "type": "doughnut",
                        "categories": [
                            "Commercial Real Estate",
                            "Commercial – General", 
                            "Commercial – Specialized",
                            "1–4 Family Residential",
                            "Auto Loans",
                            "Construction",
                            "Other Consumer"
                        ],
                        "values": [28, 27, 14, 15, 9, 4, 3],
                        "center_text": "Net Loans\n2Q'24: $2.3 Billion"
                    }
                },
                {
                    "title": "Portfolio Highlights",
                    "slide_type": "content",
                    "content": [
                        "Commercial Real Estate: Well-diversified across property types with focus on income-producing assets",
                        "Commercial – General: Strong relationships with local businesses, including SBA and PPP lending",
                        "Commercial – Specialized: Strategic focus on agricultural and energy sectors given regional expertise",
                        "1–4 Family Residential: Conservative underwriting with focus on owner-occupied properties",
                        "Auto Loans: Indirect lending through established dealer network"
                    ]
                },
                {
                    "title": "Financial Projections",
                    "slide_type": "chart", 
                    "content": ["Loan growth outlook", "Credit quality trends"],
                    "chart_data": {
                        "type": "column",
                        "categories": ["2021", "2022", "2023", "2024E"],
                        "series": [
                            {"name": "Total Loans ($B)", "values": [2.0, 2.1, 2.2, 2.3]},
                            {"name": "NPL Ratio (%)", "values": [0.8, 0.9, 0.7, 0.6]}
                        ]
                    }
                }
            ]
        }
    
    def _get_general_financial_structure(self) -> Dict[str, Any]:
        """General financial presentation structure"""
        return {
            "presentation_type": "general",
            "title": "Financial Presentation",
            "num_slides": 10,
            "design_theme": "corporate",
            "color_scheme": {"primary": "1F2937", "accent": "3B82F6"},
            "sections": [
                {
                    "title": "Title Slide",
                    "slide_type": "title",
                    "content": ["Financial Analysis", "Executive Presentation"]
                },
                {
                    "title": "Overview",
                    "slide_type": "content",
                    "content": ["Key highlights", "Financial metrics", "Strategic initiatives"]
                }
            ]
        }
    
    def generate_presentation(self, instructions: str) -> bytes:
        """Generate a complete presentation based on instructions"""
        
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not available in the Lambda environment")
        
        # Analyze the request
        structure = self.analyze_presentation_request(instructions)
        
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Generate slides based on structure
        for i, section in enumerate(structure['sections']):
            try:
                if section['slide_type'] == 'title':
                    self._create_title_slide(prs, section, structure['color_scheme'])
                elif section['slide_type'] == 'content':
                    self._create_content_slide(prs, section, structure['color_scheme'])
                elif section['slide_type'] == 'chart':
                    self._create_chart_slide(prs, section, structure['color_scheme'])
                elif section['slide_type'] == 'table':
                    self._create_table_slide(prs, section, structure['color_scheme'])
                elif section['slide_type'] == 'mixed':
                    self._create_mixed_slide(prs, section, structure['color_scheme'])
            except Exception as e:
                logger.error(f"Error creating slide {i}: {e}")
                # Continue with next slide instead of failing
                continue
        
        # Fill remaining slides if needed for PE/IB decks
        current_slides = len(prs.slides)
        target_slides = structure.get('num_slides', 10)
        
        if current_slides < target_slides:
            self._fill_remaining_slides(prs, structure, current_slides, target_slides)
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return output.read()
    
    def _create_title_slide(self, prs: Presentation, section: Dict, colors: Dict):
        """Create a professional title slide using standard layout"""
        # Use standard title slide layout
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Use existing title placeholder
        title = slide.shapes.title
        if title:
            title.text = section.get('title', 'South Plains Financial, Inc.')
            # Style the title
            for paragraph in title.text_frame.paragraphs:
                paragraph.font.size = Pt(44)
                paragraph.font.bold = True
                try:
                    paragraph.font.color.rgb = RGBColor(31, 73, 125)  # South Plains blue
                except:
                    paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Fallback to black
                paragraph.font.name = 'Arial'
                paragraph.alignment = PP_ALIGN.CENTER
        
        # Use subtitle placeholder if available
        if section.get('content') and len(slide.placeholders) > 1:
            try:
                subtitle = slide.placeholders[1]
                subtitle.text = '\n'.join(section['content'])
                for paragraph in subtitle.text_frame.paragraphs:
                    paragraph.font.size = Pt(24)
                    paragraph.font.color.rgb = RGBColor(100, 100, 100)
                    paragraph.font.name = 'Arial'
                    paragraph.alignment = PP_ALIGN.CENTER
            except:
                pass  # Skip if subtitle placeholder not available
        
        # Add simple footer text
        try:
            footer_shape = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.5))
            footer_frame = footer_shape.text_frame
            footer_frame.text = "September 2024 • CONFIDENTIAL"
            footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            footer_frame.paragraphs[0].font.size = Pt(14)
            footer_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
            footer_frame.paragraphs[0].font.name = 'Arial'
        except:
            pass  # Skip footer if error
    
    def _create_content_slide(self, prs: Presentation, section: Dict, colors: Dict):
        """Create a content slide using standard layout"""
        # Use title and content layout
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title using placeholder
        title = slide.shapes.title
        if title:
            title.text = section.get('title', 'Content')
            # Style title
            for paragraph in title.text_frame.paragraphs:
                paragraph.font.size = Pt(32)
                paragraph.font.bold = True
                try:
                    paragraph.font.color.rgb = RGBColor(31, 73, 125)  # South Plains blue
                except:
                    paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Fallback to black
                paragraph.font.name = 'Arial'
        
        # Add content using placeholder
        try:
            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            
            for i, point in enumerate(section.get('content', [])):
                if i == 0:
                    tf.text = point
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
                
                # Style each bullet point
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(55, 65, 81)  # Dark gray
                p.font.name = 'Arial'
                p.space_after = Pt(12)  # Add space between points
        except Exception as e:
            logger.error(f"Error adding content: {e}")
    
    def _create_chart_slide(self, prs: Presentation, section: Dict, colors: Dict):
        """Create a slide with charts using safe methods"""
        # Use blank layout
        slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[-1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        try:
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = section.get('title', 'Analysis')
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
            title_frame.paragraphs[0].font.name = 'Arial'
        except:
            pass
        
        # Add chart based on data
        chart_data = section.get('chart_data', {})
        
        try:
            if chart_data.get('type') == 'doughnut':
                self._create_doughnut_chart(slide, chart_data, colors)
            elif chart_data.get('type') == 'column':
                self._create_column_chart(slide, chart_data, colors)
            elif chart_data.get('type') == 'line':
                self._create_line_chart(slide, chart_data, colors)
        except Exception as e:
            logger.error(f"Error creating chart: {e}")
            # Add text instead if chart fails
            self._add_fallback_text(slide, section)
    
    def _create_doughnut_chart(self, slide, chart_data: Dict, colors: Dict):
        """Create a doughnut chart safely"""
        try:
            chart_data_obj = ChartData()
            chart_data_obj.categories = chart_data['categories']
            chart_data_obj.add_series('Portfolio', chart_data['values'])
            
            # Position chart on left side with better spacing
            x, y, cx, cy = Inches(1), Inches(2), Inches(6), Inches(4.5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data_obj
            ).chart
            
            # Style the chart
            chart.has_legend = True
            try:
                chart.legend.position = XL_LEGEND_POSITION.RIGHT
            except:
                pass  # Skip if legend position enum not available
            chart.legend.font.size = Pt(12)
            
            # Add center text if provided
            if 'center_text' in chart_data:
                center_x = x + cx/2 - Inches(1.8)
                center_y = y + cy/2 - Inches(0.4)
                textbox = slide.shapes.add_textbox(center_x, center_y, Inches(3.6), Inches(0.8))
                text_frame = textbox.text_frame
                text_frame.text = chart_data['center_text']
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                text_frame.paragraphs[0].font.bold = True
                text_frame.paragraphs[0].font.size = Pt(14)
                text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
                text_frame.paragraphs[0].font.name = 'Arial'
            
            # Add legend on the right
            self._add_custom_legend(slide, chart_data, Inches(7.5), Inches(2.5))
            
        except Exception as e:
            logger.error(f"Error in doughnut chart: {e}")
    
    def _add_custom_legend(self, slide, chart_data: Dict, x: float, y: float):
        """Add a custom legend using text boxes"""
        try:
            # Color palette for legend items
            legend_colors = [
                RGBColor(79, 129, 189),   # Theme blue
                RGBColor(192, 80, 77),    # Theme red  
                RGBColor(155, 187, 89),   # Theme green
                RGBColor(128, 100, 162),  # Theme purple
                RGBColor(247, 150, 70),   # Theme orange
                RGBColor(75, 172, 198),   # Theme teal
                RGBColor(156, 163, 175),  # Gray
            ]
            
            if chart_data.get('categories'):
                # Create legend header
                header_shape = slide.shapes.add_textbox(x, y - Inches(0.5), Inches(5), Inches(0.4))
                header_frame = header_shape.text_frame
                header_frame.text = "Portfolio Composition"
                header_frame.paragraphs[0].font.size = Pt(16)
                header_frame.paragraphs[0].font.bold = True
                header_frame.paragraphs[0].font.color.rgb = RGBColor(55, 65, 81)
                header_frame.paragraphs[0].font.name = 'Arial'
                
                for i, (category, value) in enumerate(zip(chart_data['categories'][:7], chart_data['values'][:7])):
                    y_pos = y + (i * Inches(0.5))
                    
                    # Category label with percentage
                    label_shape = slide.shapes.add_textbox(x, y_pos, Inches(4.5), Inches(0.4))
                    label_frame = label_shape.text_frame
                    label_frame.text = f"● {category}: {value}%"
                    label_frame.paragraphs[0].font.size = Pt(14)
                    label_frame.paragraphs[0].font.color.rgb = legend_colors[i % len(legend_colors)]
                    label_frame.paragraphs[0].font.name = 'Arial'
        except Exception as e:
            logger.error(f"Error adding legend: {e}")
    
    def _create_column_chart(self, slide, chart_data: Dict, colors: Dict):
        """Create a column chart safely"""
        try:
            chart_data_obj = ChartData()
            chart_data_obj.categories = chart_data['categories']
            
            for series in chart_data['series']:
                chart_data_obj.add_series(series['name'], series['values'])
            
            x, y, cx, cy = Inches(1.5), Inches(2), Inches(10), Inches(5)
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_obj
            )
        except Exception as e:
            logger.error(f"Error in column chart: {e}")
    
    def _create_line_chart(self, slide, chart_data: Dict, colors: Dict):
        """Create a line chart safely"""
        try:
            chart_data_obj = ChartData()
            chart_data_obj.categories = chart_data['categories']
            
            for series in chart_data['series']:
                chart_data_obj.add_series(series['name'], series['values'])
            
            x, y, cx, cy = Inches(1.5), Inches(2), Inches(10), Inches(5)
            slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data_obj
            )
        except Exception as e:
            logger.error(f"Error in line chart: {e}")
    
    def _create_table_slide(self, prs: Presentation, section: Dict, colors: Dict):
        """Create a slide with a table using safe methods"""
        # Use blank layout
        slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[-1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        try:
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = section.get('title', 'Table')
            title_frame.paragraphs[0].font.size = Pt(28)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
        except:
            pass
        
        # Create table
        try:
            rows = len(section.get('content', [])) + 1
            cols = 2
            left = Inches(1)
            top = Inches(2)
            width = Inches(11)
            height = Inches(0.8 * rows)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Add headers
            table.cell(0, 0).text = "Metric"
            table.cell(0, 1).text = "Value"
            
            # Add content
            for i, item in enumerate(section.get('content', [])):
                table.cell(i+1, 0).text = item
                table.cell(i+1, 1).text = "TBD"  # Would be filled with actual data
        except Exception as e:
            logger.error(f"Error creating table: {e}")
            # Add content as text if table fails
            self._add_fallback_text(slide, section)
    
    def _create_mixed_slide(self, prs: Presentation, section: Dict, colors: Dict):
        """Create a slide with mixed content using safe methods"""
        # Use title and content layout
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.title
        if title:
            title.text = section.get('title', 'Analysis')
            for paragraph in title.text_frame.paragraphs:
                paragraph.font.size = Pt(28)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(31, 73, 125)
        
        # Add content
        try:
            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            
            for i, point in enumerate(section.get('content', [])):
                if i == 0:
                    tf.text = f"• {point}"
                else:
                    p = tf.add_paragraph()
                    p.text = f"• {point}"
                    p.level = 0
        except:
            # Fallback to text boxes
            self._add_fallback_text(slide, section)
    
    def _add_fallback_text(self, slide, section: Dict):
        """Add text content when other methods fail"""
        try:
            content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(5))
            content_frame = content_shape.text_frame
            
            for i, point in enumerate(section.get('content', [])):
                if i == 0:
                    content_frame.text = f"• {point}"
                else:
                    p = content_frame.add_paragraph()
                    p.text = f"• {point}"
        except:
            pass
    
    def _fill_remaining_slides(self, prs: Presentation, structure: Dict, current: int, target: int):
        """Fill remaining slides to reach target count"""
        
        # Additional sections for PE/IB decks
        additional_sections = [
            {"title": "Financial Projections", "content": ["5-year financial model", "Revenue assumptions", "Cost structure", "Capital requirements"]},
            {"title": "Comparable Analysis", "content": ["Trading comparables", "Transaction comparables", "Valuation benchmarks"]},
            {"title": "Due Diligence Findings", "content": ["Financial due diligence", "Legal due diligence", "Operational review", "ESG considerations"]},
            {"title": "Synergy Analysis", "content": ["Revenue synergies", "Cost synergies", "Timeline to realization"]},
            {"title": "Integration Plan", "content": ["100-day plan", "Key milestones", "Integration team", "Risk mitigation"]},
            {"title": "Regulatory Considerations", "content": ["Regulatory approvals", "Compliance requirements", "Timeline"]},
            {"title": "Financing Structure", "content": ["Sources of funds", "Uses of funds", "Pro forma capitalization"]},
            {"title": "Management Incentives", "content": ["Equity participation", "Performance targets", "Retention strategy"]},
            {"title": "ESG Considerations", "content": ["Environmental impact", "Social responsibility", "Governance structure"]},
            {"title": "Appendix", "content": ["Detailed financials", "Supporting documentation", "Contact information"]}
        ]
        
        slides_needed = target - current
        for i in range(min(slides_needed, len(additional_sections))):
            section = additional_sections[i]
            section['slide_type'] = 'content'
            try:
                self._create_content_slide(prs, section, structure['color_scheme'])
            except:
                continue  # Skip if error
    
    def _apply_colors(self, slide, colors: Dict):
        """Apply color scheme to slide elements - placeholder for future enhancement"""
        pass