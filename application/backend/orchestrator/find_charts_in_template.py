#!/usr/bin/env python3
"""
Find all slides with charts in the template
"""

import boto3
import io
import logging
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(message)s')
logger = logging.getLogger(__name__)

# Initialize S3
s3 = boto3.client('s3')

def find_all_charts():
    """Find all slides with charts in the template"""
    logger.info("=" * 80)
    logger.info("SEARCHING FOR ALL CHARTS IN TEMPLATE")
    logger.info("=" * 80)
    
    # Load template from S3
    try:
        response = s3.get_object(
            Bucket='scribbe-ai-dev-documents',
            Key='PUBLIC IP South Plains (1).pptx'
        )
        template_bytes = response['Body'].read()
        logger.info(f"Loaded template: {len(template_bytes)} bytes")
    except Exception as e:
        logger.error(f"Failed to load template: {e}")
        return
    
    # Load presentation
    prs = Presentation(io.BytesIO(template_bytes))
    logger.info(f"Total slides in template: {len(prs.slides)}\n")
    
    # Search for charts
    charts_found = []
    
    for idx, slide in enumerate(prs.slides):
        slide_charts = []
        slide_title = ""
        
        # Get slide title
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text and not slide_title:
                # Get first non-empty text as title
                text = shape.text.strip()
                if text and len(text) < 100:  # Reasonable title length
                    slide_title = text
                    break
        
        # Look for charts
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'has_chart') and shape.has_chart:
                chart_info = {
                    'shape_index': shape_idx,
                    'shape_name': shape.name,
                    'chart_type': shape.chart.chart_type,
                    'has_title': shape.chart.has_title,
                }
                if shape.chart.has_title:
                    chart_info['chart_title'] = shape.chart.chart_title.text_frame.text
                
                slide_charts.append(chart_info)
        
        if slide_charts:
            charts_found.append({
                'slide_position': idx + 1,
                'slide_title': slide_title,
                'charts': slide_charts
            })
    
    # Report findings
    if charts_found:
        logger.info(f"Found {len(charts_found)} slides with charts:\n")
        
        for slide_info in charts_found:
            logger.info(f"Slide {slide_info['slide_position']}: {slide_info['slide_title']}")
            for chart in slide_info['charts']:
                logger.info(f"  - Chart: {chart['shape_name']}")
                logger.info(f"    Type: {chart['chart_type']}")
                if 'chart_title' in chart:
                    logger.info(f"    Title: {chart['chart_title']}")
            logger.info("")
    else:
        logger.warning("No charts found in any slides!")
    
    # Let's specifically look for the "Loan Portfolio" slides with visual elements
    logger.info("=" * 80)
    logger.info("ANALYZING LOAN PORTFOLIO SLIDES")
    logger.info("=" * 80)
    
    for idx, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_text.append(shape.text)
        
        full_text = ' '.join(slide_text).lower()
        
        if 'loan portfolio' in full_text:
            logger.info(f"\nSlide {idx + 1}: Loan Portfolio slide")
            
            # Count shape types
            shape_types = {}
            for shape in slide.shapes:
                shape_type = str(shape.shape_type)
                shape_types[shape_type] = shape_types.get(shape_type, 0) + 1
            
            logger.info(f"  Shape types in this slide:")
            for shape_type, count in shape_types.items():
                logger.info(f"    - {shape_type}: {count}")
            
            # Look for picture/image shapes that might be charts
            image_count = 0
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                    logger.info(f"    Found PICTURE shape: {shape.name}")
            
            if image_count > 0:
                logger.info(f"  *** This slide contains {image_count} picture(s) that might be chart images")

if __name__ == "__main__":
    find_all_charts()