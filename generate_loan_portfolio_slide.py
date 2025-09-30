#!/usr/bin/env python3
"""
Standalone script to generate the Loan Portfolio slide and upload it to S3
"""

import json
import boto3
import io
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE

# Initialize AWS S3 client
s3 = boto3.client('s3')

# S3 bucket name
OUTPUT_BUCKET = 'scribbe-ai-dev-output'

def create_loan_portfolio_slide():
    """
    Create the Loan Portfolio slide as specified
    """
    # Create a new presentation
    prs = Presentation()
    
    # Add a blank slide
    slide_layout = prs.slide_layouts[5]  # Use blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set slide dimensions (standard 16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Add title
    title_left = Inches(0.5)
    title_top = Inches(0.3)
    title_width = Inches(5)
    title_height = Inches(0.5)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = "Loan Portfolio"
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    
    # Add chart title
    chart_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(5), Inches(0.4))
    chart_title_frame = chart_title_box.text_frame
    chart_title_frame.text = "Total Loans Held for Investment ($ in Millions)"
    chart_title_frame.paragraphs[0].font.size = Pt(14)
    chart_title_frame.paragraphs[0].font.bold = True
    
    # Create combo chart (bar + line)
    chart_data = ChartData()
    chart_data.categories = ['2Q\'19', '3Q\'19', '4Q\'19', '1Q\'20', '2Q\'20']
    
    # Bar data - loan balances
    chart_data.add_series('Total Loans', (1936, 1963, 2144, 2109, 2332))
    
    # Add the chart
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(5.5), Inches(3.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    # Remove chart title (we added it separately)
    chart.has_title = False
    
    # Style the bars (red color)
    series = chart.series[0]
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(192, 0, 0)  # Red color
    
    # Add value labels on bars
    series.has_data_labels = True
    data_labels = series.data_labels
    data_labels.position = XL_CHART_TYPE.OUTSIDE_END
    data_labels.font.size = Pt(9)
    data_labels.number_format = '#,##0'
    
    # Create a second chart for the line (yield percentages)
    # Note: python-pptx doesn't support true combo charts, so we'll overlay
    line_chart_data = ChartData()
    line_chart_data.categories = ['2Q\'19', '3Q\'19', '4Q\'19', '1Q\'20', '2Q\'20']
    line_chart_data.add_series('Yield %', (5.90, 5.91, 5.79, 5.76, 5.26))
    line_chart_data.add_series('Yield with PPP', (None, None, None, None, 5.06))
    
    # Add secondary value axis labels manually (yield percentages)
    # Since python-pptx doesn't support combo charts well, we'll add text boxes for yield values
    yield_values = [(5.90, '5.90%'), (5.91, '5.91%'), (5.79, '5.79%'), (5.76, '5.76%'), (5.26, '5.26%')]
    for i, (value, label) in enumerate(yield_values):
        x_pos = Inches(0.8 + i * 1.0)
        y_pos = Inches(1.3)
        text_box = slide.shapes.add_textbox(x_pos, y_pos, Inches(0.6), Inches(0.3))
        text_frame = text_box.text_frame
        text_frame.text = label
        text_frame.paragraphs[0].font.size = Pt(10)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add PPP yield value for 2Q'20
    ppd_box = slide.shapes.add_textbox(Inches(4.8), Inches(1.6), Inches(0.6), Inches(0.3))
    ppd_frame = ppd_box.text_frame
    ppd_frame.text = "5.06%"
    ppd_frame.paragraphs[0].font.size = Pt(10)
    ppd_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    ppd_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add legend at bottom of chart
    legend_y = Inches(5.1)
    
    # Red box for Total Loans
    red_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), legend_y,
        Inches(0.2), Inches(0.15)
    )
    red_box.fill.solid()
    red_box.fill.fore_color.rgb = RGBColor(192, 0, 0)
    red_box.line.fill.background()
    
    # Legend text for Total Loans
    loans_text = slide.shapes.add_textbox(Inches(1.8), legend_y, Inches(1.5), Inches(0.2))
    loans_text.text_frame.text = "Total Loans"
    loans_text.text_frame.paragraphs[0].font.size = Pt(9)
    
    # Black line for Yield
    black_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.0), legend_y + Inches(0.05),
        Inches(0.3), Inches(0.05)
    )
    black_line.fill.solid()
    black_line.fill.fore_color.rgb = RGBColor(0, 0, 0)
    black_line.line.fill.background()
    
    # Legend text for Yield
    yield_text = slide.shapes.add_textbox(Inches(3.4), legend_y, Inches(1.5), Inches(0.2))
    yield_text.text_frame.text = "Yield"
    yield_text.text_frame.paragraphs[0].font.size = Pt(9)
    
    # Gray dashed line for Yield with PPP
    gray_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.5), legend_y + Inches(0.05),
        Inches(0.3), Inches(0.05)
    )
    gray_line.fill.solid()
    gray_line.fill.fore_color.rgb = RGBColor(128, 128, 128)
    gray_line.line.fill.background()
    
    # Legend text for Yield with PPP
    ppp_text = slide.shapes.add_textbox(Inches(4.9), legend_y, Inches(1.5), Inches(0.2))
    ppp_text.text_frame.text = "Yield with PPP"
    ppp_text.text_frame.paragraphs[0].font.size = Pt(9)
    
    # Add highlights section on the right
    highlights_left = Inches(6.3)
    highlights_top = Inches(1.2)
    highlights_width = Inches(3.5)
    highlights_height = Inches(3.8)
    
    # Highlights background (light gray)
    highlights_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        highlights_left, highlights_top,
        highlights_width, highlights_height
    )
    highlights_bg.fill.solid()
    highlights_bg.fill.fore_color.rgb = RGBColor(245, 245, 245)
    highlights_bg.line.fill.background()
    
    # Highlights title
    highlights_title = slide.shapes.add_textbox(
        highlights_left + Inches(0.2),
        highlights_top + Inches(0.1),
        highlights_width - Inches(0.4),
        Inches(0.4)
    )
    highlights_title_frame = highlights_title.text_frame
    highlights_title_frame.text = "2Q'20 Highlights"
    highlights_title_frame.paragraphs[0].font.size = Pt(16)
    highlights_title_frame.paragraphs[0].font.bold = True
    highlights_title_frame.paragraphs[0].font.color.rgb = RGBColor(192, 0, 0)  # Red accent
    
    # Highlights content
    highlights_content = slide.shapes.add_textbox(
        highlights_left + Inches(0.2),
        highlights_top + Inches(0.6),
        highlights_width - Inches(0.4),
        highlights_height - Inches(0.8)
    )
    highlights_frame = highlights_content.text_frame
    highlights_frame.word_wrap = True
    
    # Add bullet points
    highlights = [
        "Total loan increase of $229.9M vs. 1Q'20",
        "Growth from $215.3M PPP loans and $34.7M seasonal agriculture loans",
        "Partial offset from $24.4M pay-downs in non-residential consumer and direct energy loans",
        "Over 2,000 PPP loans closed",
        "2Q'20 yield of 5.26% (down 50 bps vs. 1Q'20 excluding PPP)"
    ]
    
    for i, highlight in enumerate(highlights):
        p = highlights_frame.add_paragraph() if i > 0 else highlights_frame.paragraphs[0]
        p.text = f"• {highlight}"
        p.font.size = Pt(11)
        p.space_after = Pt(8)
        p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Add company logo placeholder (top right)
    logo_text = slide.shapes.add_textbox(Inches(8.5), Inches(0.3), Inches(1.3), Inches(0.4))
    logo_frame = logo_text.text_frame
    logo_frame.text = "[Company Logo]"
    logo_frame.paragraphs[0].font.size = Pt(10)
    logo_frame.paragraphs[0].font.italic = True
    logo_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    # Add footer bar
    footer_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.2),
        Inches(10), Inches(0.425)
    )
    footer_shape.fill.solid()
    footer_shape.fill.fore_color.rgb = RGBColor(128, 128, 128)  # Gray
    footer_shape.line.fill.background()
    
    # Add footer text
    footer_text = slide.shapes.add_textbox(
        Inches(0.3), Inches(5.25),
        Inches(9.4), Inches(0.3)
    )
    footer_frame = footer_text.text_frame
    footer_frame.text = "South Plains Financial, Inc."
    footer_frame.paragraphs[0].font.size = Pt(12)
    footer_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    return prs

def upload_to_s3(presentation):
    """
    Upload the presentation to S3
    """
    # Save presentation to buffer
    output_buffer = io.BytesIO()
    presentation.save(output_buffer)
    output_buffer.seek(0)
    
    # Generate filename with timestamp
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    output_key = f"loan_portfolio_slide_{timestamp}.pptx"
    
    # Upload to S3
    s3.put_object(
        Bucket=OUTPUT_BUCKET,
        Key=output_key,
        Body=output_buffer.getvalue(),
        ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )
    
    return output_key

def main():
    """
    Main function to generate and upload the slide
    """
    print("Generating Loan Portfolio slide...")
    
    # Create the presentation
    prs = create_loan_portfolio_slide()
    
    # Upload to S3
    output_key = upload_to_s3(prs)
    
    print(f"✅ Slide successfully uploaded to S3!")
    print(f"   Bucket: {OUTPUT_BUCKET}")
    print(f"   Key: {output_key}")
    print(f"   URL: s3://{OUTPUT_BUCKET}/{output_key}")
    
    # Generate a presigned URL for downloading
    try:
        presigned_url = s3.generate_presigned_url(
            'get_object',
            Params={'Bucket': OUTPUT_BUCKET, 'Key': output_key},
            ExpiresIn=3600  # URL valid for 1 hour
        )
        print(f"\n📥 Download URL (valid for 1 hour):")
        print(f"   {presigned_url}")
    except Exception as e:
        print(f"Could not generate presigned URL: {e}")

if __name__ == "__main__":
    main()